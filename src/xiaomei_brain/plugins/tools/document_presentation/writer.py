"""Deterministic PPTX creation and common non-destructive revisions."""

from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable

from .extractor import PresentationExtractor
from xiaomei_brain.documents.presentation_project import (
    build_presentation_project,
    presentation_project_directory,
)
from .validator import validate_presentation_project


MAX_SLIDES = 200
MAX_TEXT_LENGTH = 100_000


DEFAULT_THEME = {
    "background_color": "FFFFFF",
    "title_color": "172033",
    "text_color": "354052",
    "accent_color": "4F6BED",
    "font_family": "Microsoft YaHei",
    "title_size_pt": 30,
    "body_size_pt": 18,
}


class PresentationWriter:
    format_id = "presentation"
    suffix = ".pptx"
    writer_version = "1.0.0"

    def write(
        self,
        specification: dict[str, Any],
        output_path: Path,
        *,
        source_path: Path | None = None,
        asset_paths: dict[str, Path] | None = None,
    ) -> dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValueError("演示文稿写入依赖 python-pptx 未安装") from exc

        if source_path is not None:
            if source_path.suffix.lower() != self.suffix:
                raise ValueError("Presentation writer 只能修改 PPTX 附件")
            shutil.copy2(source_path, output_path)
            presentation = Presentation(str(output_path))
            operations = specification.get("operations")
            if not isinstance(operations, list) or not operations:
                raise ValueError("修改演示文稿时 specification.operations 不能为空")
            changed = self._apply_operations(
                presentation,
                operations,
                asset_paths=asset_paths,
            )
        else:
            slides = specification.get("slides")
            if not isinstance(slides, list) or not slides:
                raise ValueError("创建演示文稿时 specification.slides 必须是非空数组")
            if len(slides) > MAX_SLIDES:
                raise ValueError(f"演示文稿不能超过 {MAX_SLIDES} 页")
            presentation = Presentation()
            self._set_page_size(presentation, specification.get("page"))
            self._set_properties(presentation, specification.get("properties"))
            theme = self._theme(specification.get("theme"))
            for slide_spec in slides:
                self._append_slide(
                    presentation,
                    slide_spec,
                    theme=theme,
                    asset_paths=asset_paths,
                )
            changed = len(slides)

        if not presentation.slides:
            raise ValueError("演示文稿至少需要一页幻灯片")
        if len(presentation.slides) > MAX_SLIDES:
            raise ValueError(f"演示文稿不能超过 {MAX_SLIDES} 页")
        presentation.save(str(output_path))

        project = build_presentation_project(
            output_path,
            presentation_project_directory(output_path),
        )
        quality = validate_presentation_project(
            presentation_project_directory(output_path),
        )

        verified = Presentation(str(output_path))
        text_count, picture_count, chart_count, note_count, characters = self._summary(verified)
        if text_count == 0 and picture_count == 0 and chart_count == 0:
            raise ValueError("生成的演示文稿没有文字或图片内容")
        extraction = PresentationExtractor().extract(output_path)
        preview = extraction.sections[0].content[:1200] if extraction.sections else ""
        return {
            "writer": self.format_id,
            "writer_version": self.writer_version,
            "validation": {
                "valid": quality["valid"],
                "delivery_ready": quality["delivery_ready"],
                "issue_count": quality["issue_count"],
                "error_count": quality["error_count"],
                "warning_count": quality["warning_count"],
                "issues": quality["issues"],
                "slide_count": len(verified.slides),
                "text_shape_count": text_count,
                "picture_count": picture_count,
                "chart_count": chart_count,
                "note_slide_count": note_count,
                "character_count": characters,
                "changed_items": changed,
                "content_preview": preview,
            },
            "presentation_project": project,
        }

    @staticmethod
    def finalize_output(temporary_path: Path, output_path: Path) -> dict[str, Any]:
        """Rebuild the companion presentation project for the final artifact path."""
        temporary_project = presentation_project_directory(temporary_path)
        if temporary_project.is_dir():
            shutil.rmtree(temporary_project)
        project = build_presentation_project(
            output_path,
            presentation_project_directory(output_path),
        )
        return {"presentation_project": project}

    @staticmethod
    def _color(value: Any, field: str):
        from pptx.dml.color import RGBColor

        text = str(value or "").strip().lstrip("#").upper()
        if len(text) != 6 or any(ch not in "0123456789ABCDEF" for ch in text):
            raise ValueError(f"{field} 必须是 6 位十六进制颜色")
        return RGBColor.from_string(text)

    @classmethod
    def _theme(cls, values: Any) -> dict[str, Any]:
        if values is not None and not isinstance(values, dict):
            raise ValueError("theme 必须是对象")
        theme = {**DEFAULT_THEME, **(values or {})}
        for key in ("background_color", "title_color", "text_color", "accent_color"):
            cls._color(theme[key], f"theme.{key}")
        for key, minimum, maximum in (
            ("title_size_pt", 12, 72),
            ("body_size_pt", 8, 48),
        ):
            size = float(theme[key])
            if not minimum <= size <= maximum:
                raise ValueError(f"theme.{key} 必须在 {minimum} 到 {maximum} 之间")
            theme[key] = size
        theme["font_family"] = str(theme["font_family"] or "").strip()
        if not theme["font_family"]:
            raise ValueError("theme.font_family 不能为空")
        return theme

    @staticmethod
    def _set_page_size(presentation: Any, values: Any) -> None:
        if values is not None and not isinstance(values, dict):
            raise ValueError("page 必须是对象")
        values = values or {}
        size = str(values.get("size") or "wide").lower()
        dimensions = {
            "wide": (33.867, 19.05),
            "standard": (25.4, 19.05),
        }
        if size not in dimensions:
            raise ValueError("page.size 仅支持 wide 或 standard")
        from pptx.util import Cm

        width, height = dimensions[size]
        presentation.slide_width = Cm(width)
        presentation.slide_height = Cm(height)

    @staticmethod
    def _set_properties(presentation: Any, values: Any) -> None:
        if not isinstance(values, dict):
            return
        properties = presentation.core_properties
        aliases = {"creator": "author", "description": "comments"}
        for key in ("title", "subject", "author", "keywords", "comments", "creator", "description"):
            if key in values:
                setattr(properties, aliases.get(key, key), str(values[key]))

    @staticmethod
    def _fill_background(slide: Any, color: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def _slide_dimensions_cm(presentation: Any) -> tuple[float, float]:
        from pptx.util import Cm

        return presentation.slide_width / Cm(1), presentation.slide_height / Cm(1)

    @classmethod
    def _append_slide(
        cls,
        presentation: Any,
        slide_spec: Any,
        *,
        theme: dict[str, Any],
        asset_paths: dict[str, Path] | None,
    ) -> Any:
        if not isinstance(slide_spec, dict):
            raise ValueError("slides 中的每一项必须是对象")
        kind = str(slide_spec.get("type") or "content").lower()
        if kind not in {"title", "section", "content", "image", "blank"}:
            raise ValueError(f"不支持的幻灯片类型: {kind}")
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide_theme = cls._theme({**theme, **(slide_spec.get("theme") or {})})
        background = slide_spec.get("background_color", slide_theme["background_color"])
        cls._fill_background(slide, cls._color(background, "slide.background_color"))
        width, height = cls._slide_dimensions_cm(presentation)

        title = str(slide_spec.get("title") or "")
        subtitle = str(slide_spec.get("subtitle") or "")
        if kind == "title":
            if title:
                cls._add_text_box(
                    slide, title, 1.5, height * 0.27, width - 3, 3.0,
                    name="XiaomeiTitle", size_pt=max(slide_theme["title_size_pt"], 36),
                    color=slide_theme["title_color"], font=slide_theme["font_family"],
                    bold=True, align="center", vertical="middle",
                )
            if subtitle:
                cls._add_text_box(
                    slide, subtitle, 3, height * 0.50, width - 6, 2.0,
                    name="XiaomeiSubtitle", size_pt=slide_theme["body_size_pt"],
                    color=slide_theme["text_color"], font=slide_theme["font_family"],
                    align="center", vertical="middle",
                )
        elif kind == "section":
            cls._add_accent(slide, width * 0.11, height * 0.34, 0.18, height * 0.30, slide_theme)
            if title:
                cls._add_text_box(
                    slide, title, width * 0.15, height * 0.34, width * 0.72, 3.0,
                    name="XiaomeiTitle", size_pt=slide_theme["title_size_pt"],
                    color=slide_theme["title_color"], font=slide_theme["font_family"], bold=True,
                )
            if subtitle:
                cls._add_text_box(
                    slide, subtitle, width * 0.15, height * 0.54, width * 0.72, 2.0,
                    name="XiaomeiSubtitle", size_pt=slide_theme["body_size_pt"],
                    color=slide_theme["text_color"], font=slide_theme["font_family"],
                )
        else:
            if title:
                cls._add_text_box(
                    slide, title, 1.4, 0.8, width - 2.8, 2.2,
                    name="XiaomeiTitle", size_pt=slide_theme["title_size_pt"],
                    color=slide_theme["title_color"], font=slide_theme["font_family"], bold=True,
                )
                cls._add_accent(slide, 1.4, 3.0, 2.2, 0.10, slide_theme)

            bullets = slide_spec.get("bullets")
            body = slide_spec.get("body")
            image = slide_spec.get("image")
            has_image = isinstance(image, dict)
            if bullets is not None or body is not None:
                body_width = width - 3.0 if not has_image else width * 0.53
                cls._add_body(
                    slide,
                    body=body,
                    bullets=bullets,
                    left=1.5,
                    top=3.6 if title else 1.5,
                    width=body_width,
                    height=height - (4.5 if title else 2.5),
                    theme=slide_theme,
                )
            if has_image:
                defaults = {
                    "x_cm": width * 0.61 if (bullets is not None or body is not None) else 2.0,
                    "y_cm": 3.6 if title else 1.5,
                    "width_cm": width * 0.34 if (bullets is not None or body is not None) else width - 4.0,
                    "height_cm": height - (4.8 if title else 3.0),
                }
                cls._add_image(slide, {**defaults, **image}, asset_paths)

        elements = slide_spec.get("elements", [])
        if not isinstance(elements, list):
            raise ValueError("slide.elements 必须是数组")
        for element in elements:
            cls._add_element(slide, element, slide_theme, asset_paths)
        cls._set_notes(slide, slide_spec.get("notes"))
        return slide

    @classmethod
    def _add_text_box(
        cls,
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        name: str,
        size_pt: float,
        color: str,
        font: str,
        bold: bool = False,
        align: str = "left",
        vertical: str = "top",
    ) -> Any:
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Cm, Pt

        cls._validate_box(left, top, width, height)
        shape = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        shape.name = name
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(vertical, MSO_ANCHOR.TOP)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = cls._text(text)
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = cls._color(color, "text.color")
        return shape

    @classmethod
    def _add_body(
        cls,
        slide: Any,
        *,
        body: Any,
        bullets: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        theme: dict[str, Any],
    ) -> None:
        if bullets is not None and not isinstance(bullets, list):
            raise ValueError("slide.bullets 必须是数组")
        lines: list[tuple[str, int, bool]] = []
        if body is not None:
            lines.append((str(body), 0, False))
        for item in bullets or []:
            if isinstance(item, dict):
                text = str(item.get("text") or "")
                level = int(item.get("level", 0))
            else:
                text = str(item)
                level = 0
            if not 0 <= level <= 5:
                raise ValueError("bullet.level 必须在 0 到 5 之间")
            lines.append((text, level, True))
        if not lines:
            return
        shape = cls._add_text_box(
            slide, "", left, top, width, height,
            name="XiaomeiBody", size_pt=theme["body_size_pt"],
            color=theme["text_color"], font=theme["font_family"],
        )
        frame = shape.text_frame
        frame.clear()
        from pptx.util import Pt

        for index, (text, level, bullet) in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.level = level
            paragraph.space_after = Pt(8)
            run = paragraph.add_run()
            run.text = ("• " if bullet else "") + cls._text(text)
            run.font.name = theme["font_family"]
            run.font.size = Pt(theme["body_size_pt"] - min(level, 2))
            run.font.color.rgb = cls._color(theme["text_color"], "theme.text_color")

    @classmethod
    def _add_accent(
        cls,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        theme: dict[str, Any],
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Cm

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(left), Cm(top), Cm(width), Cm(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = cls._color(theme["accent_color"], "theme.accent_color")
        shape.line.fill.background()

    @classmethod
    def _add_image(
        cls,
        slide: Any,
        values: Any,
        asset_paths: dict[str, Path] | None,
    ) -> Any:
        if not isinstance(values, dict):
            raise ValueError("image 必须是对象")
        attachment_id = str(values.get("attachment_id") or "").strip()
        workspace_path = str(values.get("workspace_path") or "").strip()
        if bool(attachment_id) == bool(workspace_path):
            raise ValueError("image 必须且只能提供 attachment_id 或 workspace_path 之一")
        key = attachment_id or f"workspace:{workspace_path}"
        image_path = (asset_paths or {}).get(key)
        if image_path is None or not image_path.is_file():
            raise ValueError(f"当前执行现场没有可用图片: {attachment_id or workspace_path}")
        left = float(values.get("x_cm", 1.5))
        top = float(values.get("y_cm", 1.5))
        width = float(values.get("width_cm", 12))
        height_value = values.get("height_cm")
        height = float(height_value) if height_value is not None else None
        cls._validate_box(left, top, width, height or 1)
        from pptx.util import Cm

        kwargs = {"width": Cm(width)}
        if height is not None:
            kwargs["height"] = Cm(height)
        return slide.shapes.add_picture(str(image_path), Cm(left), Cm(top), **kwargs)

    @classmethod
    def _add_element(
        cls,
        slide: Any,
        element: Any,
        theme: dict[str, Any],
        asset_paths: dict[str, Path] | None,
    ) -> None:
        if not isinstance(element, dict):
            raise ValueError("slide.elements 中的每一项必须是对象")
        kind = str(element.get("type") or "text").lower()
        if kind == "image":
            cls._add_image(slide, element, asset_paths)
            return
        if kind == "shape":
            cls._add_shape_element(slide, element, theme)
            return
        if kind == "line":
            cls._add_line_element(slide, element, theme)
            return
        if kind == "table":
            cls._add_table_element(slide, element, theme)
            return
        if kind == "chart":
            cls._add_chart_element(slide, element, theme)
            return
        if kind != "text":
            raise ValueError(
                f"不支持的 slide element: {kind}；支持 text、image、shape、line、table、chart"
            )
        required = ("x_cm", "y_cm", "width_cm", "height_cm")
        if any(key not in element for key in required):
            raise ValueError("text element 必须提供 x_cm、y_cm、width_cm 和 height_cm")
        cls._add_text_box(
            slide,
            str(element.get("text") or ""),
            float(element["x_cm"]),
            float(element["y_cm"]),
            float(element["width_cm"]),
            float(element["height_cm"]),
            name=str(element.get("name") or "XiaomeiText"),
            size_pt=float(element.get("size_pt", theme["body_size_pt"])),
            color=str(element.get("color") or theme["text_color"]),
            font=str(element.get("font") or theme["font_family"]),
            bold=element.get("bold") is True,
            align=str(element.get("align") or "left").lower(),
            vertical=str(element.get("vertical") or "top").lower(),
        )

    @classmethod
    def _element_box(cls, values: dict[str, Any], kind: str) -> tuple[float, float, float, float]:
        required = ("x_cm", "y_cm", "width_cm", "height_cm")
        if any(key not in values for key in required):
            raise ValueError(
                f"{kind} element 必须提供 x_cm、y_cm、width_cm 和 height_cm"
            )
        try:
            box = tuple(float(values[key]) for key in required)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"{kind} element 的坐标和尺寸必须是数字") from exc
        cls._validate_box(*box)
        return box

    @classmethod
    def _format_text_frame(
        cls,
        frame: Any,
        text: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
    ) -> None:
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Pt

        size = float(values.get("font_size_pt", values.get("size_pt", theme["body_size_pt"])))
        if not 6 <= size <= 144:
            raise ValueError("font_size_pt 必须在 6 到 144 之间")
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(str(values.get("vertical") or "middle").lower(), MSO_ANCHOR.MIDDLE)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(str(values.get("align") or "center").lower(), PP_ALIGN.CENTER)
        run = paragraph.add_run()
        run.text = cls._text(text)
        run.font.name = str(values.get("font") or theme["font_family"])
        run.font.size = Pt(size)
        run.font.bold = values.get("bold") is True
        run.font.color.rgb = cls._color(
            values.get("text_color") or values.get("color") or theme["text_color"],
            "text_color",
        )

    @classmethod
    def _apply_line_style(
        cls,
        shape: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
        *,
        default_color: str | None,
    ) -> None:
        from pptx.util import Pt

        line_color = values.get("line_color", default_color)
        if line_color is None or str(line_color).strip().lower() in {"none", "transparent"}:
            shape.line.fill.background()
            return
        shape.line.color.rgb = cls._color(line_color, "line_color")
        width = float(values.get("line_width_pt", 1.5))
        if not 0.1 <= width <= 50:
            raise ValueError("line_width_pt 必须在 0.1 到 50 之间")
        shape.line.width = Pt(width)
        if "line_dash" in values:
            cls._set_line_dash(shape, values["line_dash"])
        for arrow_field in ("start_arrow", "end_arrow"):
            if arrow_field in values:
                cls._set_line_arrow(shape, arrow_field, values[arrow_field])
        if "line_transparency" in values:
            cls._set_solid_fill_transparency(
                shape, values["line_transparency"], line=True,
            )

    @classmethod
    def _add_shape_element(
        cls,
        slide: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
    ) -> Any:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Cm

        aliases = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rect": MSO_SHAPE.RECTANGLE,
            "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "hexagon": MSO_SHAPE.HEXAGON,
            "chevron": MSO_SHAPE.CHEVRON,
            "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "trapezoid": MSO_SHAPE.TRAPEZOID,
        }
        shape_kind = str(values.get("shape") or "rectangle").strip().lower().replace("-", "_")
        auto_shape = aliases.get(shape_kind)
        if auto_shape is None:
            raise ValueError(
                "shape 仅支持 rectangle、round_rect、ellipse、triangle、diamond、"
                "hexagon、chevron、pentagon、parallelogram、trapezoid"
            )
        left, top, width, height = cls._element_box(values, "shape")
        shape = slide.shapes.add_shape(
            auto_shape, Cm(left), Cm(top), Cm(width), Cm(height),
        )
        shape.name = str(values.get("name") or "XiaomeiShape")
        fill_color = values.get("fill_color", theme["accent_color"])
        if fill_color is None or str(fill_color).strip().lower() in {"none", "transparent"}:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = cls._color(fill_color, "fill_color")
            if "fill_transparency" in values:
                cls._set_solid_fill_transparency(
                    shape, values["fill_transparency"], line=False,
                )
        cls._apply_line_style(
            shape, values, theme, default_color=None,
        )
        if "text" in values:
            cls._format_text_frame(shape.text_frame, values["text"], values, theme)
        if "rotation" in values:
            rotation = float(values["rotation"])
            if not -360 <= rotation <= 360:
                raise ValueError("rotation 必须在 -360 到 360 之间")
            shape.rotation = rotation
        return shape

    @classmethod
    def _add_line_element(
        cls,
        slide: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
    ) -> Any:
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.util import Cm

        required = ("x_cm", "y_cm", "to_x_cm", "to_y_cm")
        if any(key not in values for key in required):
            raise ValueError("line element 必须提供 x_cm、y_cm、to_x_cm 和 to_y_cm")
        try:
            start_x, start_y, end_x, end_y = (
                float(values[key]) for key in required
            )
        except (TypeError, ValueError) as exc:
            raise ValueError("line element 的起点和终点必须是数字") from exc
        if min(start_x, start_y, end_x, end_y) < 0 or max(start_x, start_y, end_x, end_y) > 200:
            raise ValueError("line element 的坐标必须在 0 到 200 厘米之间")
        connectors = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curve": MSO_CONNECTOR.CURVE,
        }
        connector_name = str(values.get("connector") or "straight").strip().lower()
        connector = connectors.get(connector_name)
        if connector is None:
            raise ValueError("connector 仅支持 straight、elbow、curve")
        shape = slide.shapes.add_connector(
            connector, Cm(start_x), Cm(start_y), Cm(end_x), Cm(end_y),
        )
        shape.name = str(values.get("name") or "XiaomeiLine")
        cls._apply_line_style(
            shape, values, theme, default_color=theme["accent_color"],
        )
        return shape

    @classmethod
    def _format_table_cell(
        cls,
        cell: Any,
        value: Any,
        defaults: dict[str, Any],
        theme: dict[str, Any],
    ) -> None:
        values = dict(defaults)
        if isinstance(value, dict):
            values.update(value)
            text = value.get("text", "")
        else:
            text = value
        if "fill_color" in values:
            fill_color = values.get("fill_color")
            if fill_color is None or str(fill_color).strip().lower() in {"none", "transparent"}:
                cell.fill.background()
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cls._color(fill_color, "table.fill_color")
        values.setdefault("align", "left")
        values.setdefault("vertical", "middle")
        cls._format_text_frame(cell.text_frame, text, values, theme)

    @classmethod
    def _add_table_element(
        cls,
        slide: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
    ) -> Any:
        from pptx.util import Cm

        left, top, width, height = cls._element_box(values, "table")
        data = values.get("data")
        if not isinstance(data, list) or not data:
            raise ValueError("table.data 必须是非空二维数组")
        if len(data) > 200:
            raise ValueError("table.data 不能超过 200 行")
        if not isinstance(data[0], list) or not data[0]:
            raise ValueError("table.data 每一行必须是非空数组")
        column_count = len(data[0])
        if column_count > 50:
            raise ValueError("table.data 不能超过 50 列")
        if any(not isinstance(row, list) or len(row) != column_count for row in data):
            raise ValueError("table.data 的每一行必须具有相同列数")
        shape = slide.shapes.add_table(
            len(data), column_count, Cm(left), Cm(top), Cm(width), Cm(height),
        )
        shape.name = str(values.get("name") or "XiaomeiTable")
        table = shape.table
        widths = values.get("column_widths_cm")
        if widths is not None:
            if not isinstance(widths, list) or len(widths) != column_count:
                raise ValueError("column_widths_cm 必须与 table.data 列数相同")
            normalized_widths = [float(item) for item in widths]
            if any(item <= 0 for item in normalized_widths):
                raise ValueError("column_widths_cm 中的宽度必须大于 0")
            for index, column_width in enumerate(normalized_widths):
                table.columns[index].width = Cm(column_width)
        cell_defaults = values.get("cell_style") or {}
        header_defaults = values.get("header_style") or {}
        if not isinstance(cell_defaults, dict) or not isinstance(header_defaults, dict):
            raise ValueError("cell_style 和 header_style 必须是对象")
        for row_index, row in enumerate(data):
            for column_index, value in enumerate(row):
                defaults = cell_defaults
                if row_index == 0 and header_defaults:
                    defaults = {**cell_defaults, **header_defaults}
                cls._format_table_cell(
                    table.cell(row_index, column_index), value, defaults, theme,
                )
        return shape

    @classmethod
    def _chart_data(cls, values: dict[str, Any]) -> Any:
        from pptx.chart.data import CategoryChartData

        categories = values.get("categories")
        series_items = values.get("series")
        if not isinstance(categories, list) or not categories:
            raise ValueError("chart.categories 必须是非空数组")
        if not isinstance(series_items, list) or not series_items:
            raise ValueError("chart.series 必须是非空数组")
        if len(categories) > 2_000 or len(series_items) > 100:
            raise ValueError("图表数据超过支持的大小")
        data = CategoryChartData()
        data.categories = [str(value) for value in categories]
        for index, item in enumerate(series_items):
            if not isinstance(item, dict):
                raise ValueError("chart.series 中的每一项必须是对象")
            series_values = item.get("values")
            if not isinstance(series_values, list) or len(series_values) != len(categories):
                raise ValueError("每个 chart.series 必须为每个 category 提供一个值")
            normalized: list[float | None] = []
            for value in series_values:
                if value is None:
                    normalized.append(None)
                    continue
                try:
                    normalized.append(float(value))
                except (TypeError, ValueError) as exc:
                    raise ValueError(f"无效的图表数值: {value}") from exc
            data.add_series(str(item.get("name") or f"系列 {index + 1}"), normalized)
        return data

    @classmethod
    def _add_chart_element(
        cls,
        slide: Any,
        values: dict[str, Any],
        theme: dict[str, Any],
    ) -> Any:
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Cm

        left, top, width, height = cls._element_box(values, "chart")
        chart_types = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
        }
        type_name = str(values.get("chart_type") or "column").strip().lower().replace("-", "_")
        chart_type = chart_types.get(type_name)
        if chart_type is None:
            raise ValueError(
                "chart_type 仅支持 column、column_stacked、bar、line、line_markers、"
                "pie、doughnut、area"
            )
        shape = slide.shapes.add_chart(
            chart_type, Cm(left), Cm(top), Cm(width), Cm(height), cls._chart_data(values),
        )
        shape.name = str(values.get("name") or "XiaomeiChart")
        chart = shape.chart
        title = cls._text(values.get("title"))
        chart.has_title = bool(title)
        if title:
            chart.chart_title.text_frame.text = title
        chart.has_legend = bool(values.get("show_legend", True))
        if chart.has_legend:
            positions = {
                "top": XL_LEGEND_POSITION.TOP,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
            }
            position_name = str(values.get("legend_position") or "bottom").lower()
            if position_name not in positions:
                raise ValueError("legend_position 仅支持 top、bottom、left、right")
            chart.legend.position = positions[position_name]
        colors = values.get("series_colors")
        if colors is not None:
            if not isinstance(colors, list) or not colors:
                raise ValueError("series_colors 必须是非空数组")
            for index, series in enumerate(chart.series):
                rgb = cls._color(colors[index % len(colors)], f"series_colors[{index}]")
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = rgb
                except (AttributeError, ValueError):
                    pass
                try:
                    series.format.line.color.rgb = rgb
                except (AttributeError, ValueError):
                    pass
        if values.get("show_values") is True:
            for plot in chart.plots:
                plot.has_data_labels = True
                plot.data_labels.show_value = True
        return shape

    @staticmethod
    def _validate_box(left: float, top: float, width: float, height: float) -> None:
        if left < 0 or top < 0 or width <= 0 or height <= 0:
            raise ValueError("元素坐标必须非负，宽高必须大于 0")
        if max(left, top, width, height) > 200:
            raise ValueError("元素坐标或尺寸超过 200 厘米限制")

    @staticmethod
    def _text(value: Any) -> str:
        text = str(value or "")
        if len(text) > MAX_TEXT_LENGTH:
            raise ValueError(f"单个文本内容不能超过 {MAX_TEXT_LENGTH} 个字符")
        return text

    @classmethod
    def _set_notes(cls, slide: Any, value: Any) -> None:
        if value is None:
            return
        notes = cls._text(value)
        slide.notes_slide.notes_text_frame.text = notes

    @staticmethod
    def _shape_named(slide: Any, name: str) -> Any | None:
        return next((shape for shape in slide.shapes if shape.name == name), None)

    @classmethod
    def _set_shape_text(cls, shape: Any, text: Any) -> None:
        if shape is None or not getattr(shape, "has_text_frame", False):
            raise ValueError("幻灯片中没有可更新的目标文本框")
        shape.text_frame.text = cls._text(text)

    @classmethod
    def _replace_in_frame(cls, frame: Any, old: str, new: str, replace_all: bool) -> int:
        replacements = 0
        for paragraph in frame.paragraphs:
            runs = list(paragraph.runs)
            if not runs:
                continue
            content = "".join(run.text for run in runs)
            available = content.count(old)
            if not available:
                continue
            count = available if replace_all else 1
            updated = content.replace(old, new, count)
            runs[0].text = updated
            for run in runs[1:]:
                run.text = ""
            replacements += count
            if not replace_all:
                return replacements
        return replacements

    @classmethod
    def _replace_text(
        cls,
        presentation: Any,
        old: str,
        new: str,
        *,
        replace_all: bool,
        required: bool = True,
    ) -> int:
        replacements = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    count = cls._replace_in_frame(shape.text_frame, old, new, replace_all)
                    replacements += count
                    if count and not replace_all:
                        return replacements
            try:
                if not slide.has_notes_slide:
                    continue
                count = cls._replace_in_frame(
                    slide.notes_slide.notes_text_frame,
                    old,
                    new,
                    replace_all,
                )
            except (AttributeError, ValueError):
                count = 0
            replacements += count
            if count and not replace_all:
                return replacements
        if required and replacements == 0:
            raise ValueError(f"演示文稿中没有找到文本: {old}")
        return replacements

    @classmethod
    def _update_slide(
        cls,
        presentation: Any,
        operation: dict[str, Any],
        asset_paths: dict[str, Path] | None,
    ) -> int:
        index = cls._slide_index(presentation, operation.get("slide"))
        slide = presentation.slides[index]
        changed = 0
        for key, shape_name in (
            ("title", "XiaomeiTitle"),
            ("subtitle", "XiaomeiSubtitle"),
            ("body", "XiaomeiBody"),
        ):
            if key not in operation:
                continue
            shape = cls._shape_named(slide, shape_name)
            if shape is None and key == "title":
                shape = slide.shapes.title
            cls._set_shape_text(shape, operation[key])
            changed += 1
        if "notes" in operation:
            cls._set_notes(slide, operation["notes"])
            changed += 1
        if "background_color" in operation:
            cls._fill_background(
                slide,
                cls._color(operation["background_color"], "background_color"),
            )
            changed += 1
        elements = operation.get("elements", [])
        if not isinstance(elements, list):
            raise ValueError("update_slide.elements 必须是数组")
        theme = cls._theme(operation.get("theme"))
        for element in elements:
            cls._add_element(slide, element, theme, asset_paths)
            changed += 1
        return changed

    @staticmethod
    def _slide_index(presentation: Any, value: Any) -> int:
        try:
            number = int(value)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"无效幻灯片页码: {value}") from exc
        if not 1 <= number <= len(presentation.slides):
            raise ValueError(f"幻灯片页码超出范围: {number}")
        return number - 1

    @staticmethod
    def _delete_slide(presentation: Any, index: int) -> None:
        slide_id = presentation.slides._sldIdLst[index]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[index]

    @staticmethod
    def _move_slide(presentation: Any, source: int, target: int) -> None:
        slide_ids = presentation.slides._sldIdLst
        element = slide_ids[source]
        slide_ids.remove(element)
        slide_ids.insert(target, element)

    @staticmethod
    def _shape_path(element_id: Any, slide_number: int) -> tuple[str, list[int]]:
        stable_prefix = f"slide-{slide_number}-shape-id-"
        prefix = f"slide-{slide_number}-shape-"
        value = str(element_id or "").strip()
        if value.startswith(stable_prefix):
            try:
                path = [int(part) for part in value[len(stable_prefix):].split(".")]
            except ValueError as exc:
                raise ValueError(f"Invalid presentation element ID: {value}") from exc
            if not path or any(shape_id < 1 for shape_id in path):
                raise ValueError(f"Invalid presentation element ID: {value}")
            return "shape_id", path
        if not value.startswith(prefix):
            raise ValueError(f"无效的演示元素 ID: {value}")
        encoded = value[len(prefix):]
        if "." in encoded:
            try:
                path = [int(part) for part in encoded.split(".")]
            except ValueError as exc:
                raise ValueError(f"无效的演示元素 ID: {value}") from exc
        else:
            try:
                number = int(encoded)
            except ValueError as exc:
                raise ValueError(f"无效的演示元素 ID: {value}") from exc
            path = []
            while number > 999:
                path.insert(0, number % 1000)
                number //= 1000
            path.insert(0, number)
        if not path or any(index < 1 for index in path):
            raise ValueError(f"无效的演示元素 ID: {value}")
        return "position", path

    @classmethod
    def _resolve_shape(cls, presentation: Any, operation: dict[str, Any]) -> Any:
        slide_index = cls._slide_index(presentation, operation.get("slide"))
        slide_number = slide_index + 1
        path_kind, path = cls._shape_path(operation.get("element_id"), slide_number)
        shapes = presentation.slides[slide_index].shapes
        shape = None
        for depth, identifier in enumerate(path):
            if path_kind == "position" and identifier > len(shapes):
                raise ValueError(f"演示元素不存在: {operation.get('element_id')}")
            if path_kind == "shape_id":
                shape = next(
                    (candidate for candidate in shapes if int(candidate.shape_id) == identifier),
                    None,
                )
                if shape is None:
                    raise ValueError(
                        f"Presentation element does not exist: {operation.get('element_id')}"
                    )
            else:
                shape = shapes[identifier - 1]
            if depth < len(path) - 1:
                nested = getattr(shape, "shapes", None)
                if nested is None:
                    raise ValueError(f"演示元素路径无效: {operation.get('element_id')}")
                shapes = nested
        return shape

    @staticmethod
    def _line_element(shape: Any) -> Any:
        try:
            return shape._element.spPr.get_or_add_ln()
        except (AttributeError, TypeError) as exc:
            raise ValueError("选中的演示元素不支持线条样式") from exc

    @staticmethod
    def _replace_xml_child(parent: Any, tag: str, child: Any) -> None:
        from pptx.oxml.ns import qn

        existing = parent.find(qn(tag))
        if existing is not None:
            parent.remove(existing)
        if child is not None:
            parent.append(child)
        order = {
            "noFill": 10, "solidFill": 10, "gradFill": 10, "pattFill": 10,
            "prstDash": 20, "custDash": 20,
            "round": 30, "bevel": 30, "miter": 30,
            "headEnd": 40, "tailEnd": 50, "extLst": 60,
        }
        children = list(parent)
        children.sort(key=lambda item: order.get(str(item.tag).rsplit("}", 1)[-1], 35))
        for item in children:
            parent.remove(item)
            parent.append(item)

    @classmethod
    def _set_line_dash(cls, shape: Any, value: Any) -> None:
        from pptx.oxml.xmlchemy import OxmlElement

        normalized = str(value or "").strip().lower().replace("-", "_")
        aliases = {
            "solid": "solid",
            "dash": "dash",
            "dashed": "dash",
            "dot": "dot",
            "dotted": "dot",
            "dash_dot": "dashDot",
            "long_dash": "lgDash",
            "long_dash_dot": "lgDashDot",
        }
        preset = aliases.get(normalized)
        if preset is None:
            raise ValueError(
                "line_dash 仅支持 solid、dash、dot、dash_dot、long_dash、long_dash_dot"
            )
        node = OxmlElement("a:prstDash")
        node.set("val", preset)
        cls._replace_xml_child(cls._line_element(shape), "a:prstDash", node)

    @classmethod
    def _set_line_arrow(cls, shape: Any, tag: str, value: Any) -> None:
        from pptx.oxml.xmlchemy import OxmlElement

        values = value if isinstance(value, dict) else {"type": value}
        arrow_type = str(values.get("type") or "none").strip().lower()
        aliases = {"arrow": "triangle", "circle": "oval"}
        arrow_type = aliases.get(arrow_type, arrow_type)
        if arrow_type not in {"none", "triangle", "stealth", "diamond", "oval", "open"}:
            raise ValueError(
                f"{tag} 仅支持 none、triangle、stealth、diamond、oval、open"
            )
        width = str(values.get("width") or "med").strip().lower()
        length = str(values.get("length") or "med").strip().lower()
        if width not in {"sm", "med", "lg"} or length not in {"sm", "med", "lg"}:
            raise ValueError(f"{tag} 的 width 和 length 仅支持 sm、med、lg")
        node = OxmlElement("a:headEnd" if tag == "start_arrow" else "a:tailEnd")
        node.set("type", arrow_type)
        node.set("w", width)
        node.set("len", length)
        cls._replace_xml_child(
            cls._line_element(shape),
            "a:headEnd" if tag == "start_arrow" else "a:tailEnd",
            node,
        )

    @staticmethod
    def _transparency(value: Any, field: str) -> float:
        try:
            number = float(value)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"{field} 必须是 0 到 100 之间的数字") from exc
        if not 0 <= number <= 100:
            raise ValueError(f"{field} 必须在 0 到 100 之间")
        return number

    @classmethod
    def _set_solid_fill_transparency(
        cls,
        shape: Any,
        value: Any,
        *,
        line: bool,
    ) -> None:
        from pptx.oxml.ns import qn
        from pptx.oxml.xmlchemy import OxmlElement

        transparency = cls._transparency(
            value,
            "line_transparency" if line else "fill_transparency",
        )
        properties = shape._element.spPr
        container = cls._line_element(shape) if line else properties
        solid = container.find(qn("a:solidFill"))
        if solid is None and line:
            style = shape._element.find(qn("p:style"))
            reference = style.find(qn("a:lnRef")) if style is not None else None
            color = next(iter(reference), None) if reference is not None else None
            solid = OxmlElement("a:solidFill")
            if color is not None:
                solid.append(deepcopy(color))
            else:
                fallback = OxmlElement("a:schemeClr")
                fallback.set("val", "accent1")
                solid.append(fallback)
            container.insert(0, solid)
        if solid is None:
            raise ValueError(
                "fill_transparency 仅支持纯色填充；可在同一操作中同时提供 fill_color"
            )
        color = next(iter(solid), None)
        if color is None:
            raise ValueError("无法修改当前颜色的透明度")
        alpha = color.find(qn("a:alpha"))
        if alpha is None:
            alpha = OxmlElement("a:alpha")
            color.append(alpha)
        alpha.set("val", str(round((100 - transparency) * 1000)))

    @classmethod
    def _update_element(cls, presentation: Any, operation: dict[str, Any]) -> int:
        from pptx.util import Cm, Pt

        shape = cls._resolve_shape(presentation, operation)
        changed = 0
        if "text" in operation:
            if not getattr(shape, "has_text_frame", False):
                raise ValueError("选中的演示元素不支持修改文字")
            frame = shape.text_frame
            paragraphs = list(frame.paragraphs)
            runs = [run for paragraph in paragraphs for run in paragraph.runs]
            value = cls._text(operation["text"])
            if runs:
                runs[0].text = value
                for run in runs[1:]:
                    run.text = ""
                for paragraph in paragraphs[1:]:
                    paragraph.text = ""
            else:
                frame.text = value
            changed += 1
        if "fill_color" in operation:
            shape.fill.solid()
            shape.fill.fore_color.rgb = cls._color(operation["fill_color"], "fill_color")
            changed += 1
        if "line_color" in operation:
            shape.line.color.rgb = cls._color(operation["line_color"], "line_color")
            changed += 1
        if "line_width_pt" in operation:
            width = float(operation["line_width_pt"])
            if not 0.1 <= width <= 50:
                raise ValueError("line_width_pt 必须在 0.1 到 50 之间")
            shape.line.width = Pt(width)
            changed += 1
        if "line_dash" in operation:
            cls._set_line_dash(shape, operation["line_dash"])
            changed += 1
        for arrow_field in ("start_arrow", "end_arrow"):
            if arrow_field in operation:
                cls._set_line_arrow(shape, arrow_field, operation[arrow_field])
                changed += 1
        if "fill_transparency" in operation:
            cls._set_solid_fill_transparency(
                shape, operation["fill_transparency"], line=False,
            )
            changed += 1
        if "line_transparency" in operation:
            cls._set_solid_fill_transparency(
                shape, operation["line_transparency"], line=True,
            )
            changed += 1
        text_style_fields = {"text_color", "font_size_pt", "bold"}
        if text_style_fields.intersection(operation):
            if not getattr(shape, "has_text_frame", False):
                raise ValueError("选中的演示元素不支持修改文字样式")
            runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
            if not runs:
                runs = [shape.text_frame.paragraphs[0].add_run()]
            if "font_size_pt" in operation:
                size = float(operation["font_size_pt"])
                if not 6 <= size <= 144:
                    raise ValueError("font_size_pt 必须在 6 到 144 之间")
            for run in runs:
                if "text_color" in operation:
                    run.font.color.rgb = cls._color(operation["text_color"], "text_color")
                if "font_size_pt" in operation:
                    run.font.size = Pt(float(operation["font_size_pt"]))
                if "bold" in operation:
                    run.font.bold = bool(operation["bold"])
            changed += len(text_style_fields.intersection(operation))
        for field, attribute in (
            ("x_cm", "left"),
            ("y_cm", "top"),
            ("width_cm", "width"),
            ("height_cm", "height"),
        ):
            if field not in operation:
                continue
            value = float(operation[field])
            if value < 0 or (field in {"width_cm", "height_cm"} and value <= 0):
                raise ValueError(f"{field} 的值无效")
            setattr(shape, attribute, Cm(value))
            changed += 1
        if changed == 0:
            raise ValueError("update_element 没有提供可修改的字段")
        return changed

    @classmethod
    def _delete_element(cls, presentation: Any, operation: dict[str, Any]) -> None:
        shape = cls._resolve_shape(presentation, operation)
        element = shape._element
        parent = element.getparent()
        if parent is None:
            raise ValueError(f"演示元素无法删除: {operation.get('element_id')}")
        parent.remove(element)

    @classmethod
    def _update_table_cell(cls, presentation: Any, operation: dict[str, Any]) -> int:
        from pptx.util import Pt

        shape = cls._resolve_shape(presentation, operation)
        if not getattr(shape, "has_table", False):
            raise ValueError("The selected presentation element is not a table")
        try:
            row = int(operation.get("row"))
            column = int(operation.get("column"))
        except (TypeError, ValueError) as exc:
            raise ValueError(
                "update_table_cell requires 1-based row and column values"
            ) from exc
        table = shape.table
        if not 1 <= row <= len(table.rows) or not 1 <= column <= len(table.columns):
            raise ValueError(f"Table cell is outside the valid range: {row},{column}")
        cell = table.cell(row - 1, column - 1)
        changed = 0
        if "text" in operation:
            cell.text = cls._text(operation["text"])
            changed += 1
        if "fill_color" in operation:
            cell.fill.solid()
            cell.fill.fore_color.rgb = cls._color(operation["fill_color"], "fill_color")
            changed += 1
        style_fields = {"text_color", "font_size_pt", "bold"}
        if style_fields.intersection(operation):
            if "font_size_pt" in operation:
                size = float(operation["font_size_pt"])
                if not 6 <= size <= 144:
                    raise ValueError("font_size_pt must be between 6 and 144")
            runs = [run for paragraph in cell.text_frame.paragraphs for run in paragraph.runs]
            if not runs:
                runs = [cell.text_frame.paragraphs[0].add_run()]
            for run in runs:
                if "text_color" in operation:
                    run.font.color.rgb = cls._color(operation["text_color"], "text_color")
                if "font_size_pt" in operation:
                    run.font.size = Pt(float(operation["font_size_pt"]))
                if "bold" in operation:
                    run.font.bold = bool(operation["bold"])
            changed += len(style_fields.intersection(operation))
        if changed == 0:
            raise ValueError("update_table_cell did not include any editable fields")
        return changed

    @classmethod
    def _replace_image(
        cls,
        presentation: Any,
        operation: dict[str, Any],
        asset_paths: dict[str, Path] | None,
    ) -> int:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        shape = cls._resolve_shape(presentation, operation)
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError("The selected presentation element is not an image")
        attachment_id = str(operation.get("attachment_id") or "").strip()
        workspace_path = str(operation.get("workspace_path") or "").strip()
        if bool(attachment_id) == bool(workspace_path):
            raise ValueError(
                "replace_image requires exactly one of attachment_id or workspace_path"
            )
        key = attachment_id or f"workspace:{workspace_path}"
        image_path = (asset_paths or {}).get(key)
        if image_path is None or not image_path.is_file():
            raise ValueError(
                f"The replacement image is unavailable: {attachment_id or workspace_path}"
            )
        _image_part, relationship_id = shape.part.get_or_add_image_part(str(image_path))
        shape._element.blipFill.blip.rEmbed = relationship_id
        for field in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
            setattr(shape, field, 0)
        return 1

    @classmethod
    def _update_chart(cls, presentation: Any, operation: dict[str, Any]) -> int:
        shape = cls._resolve_shape(presentation, operation)
        if not getattr(shape, "has_chart", False):
            raise ValueError("The selected presentation element is not a chart")
        chart = shape.chart
        changed = 0
        if "title" in operation:
            title = cls._text(operation.get("title"))
            chart.has_title = bool(title)
            if title:
                chart.chart_title.text_frame.text = title
            changed += 1
        categories_supplied = "categories" in operation
        series_supplied = "series" in operation
        if categories_supplied != series_supplied:
            raise ValueError("update_chart must provide categories and series together")
        if categories_supplied:
            from pptx.chart.data import CategoryChartData

            categories = operation.get("categories")
            series_items = operation.get("series")
            if not isinstance(categories, list) or not categories:
                raise ValueError("update_chart.categories must be a non-empty array")
            if not isinstance(series_items, list) or not series_items:
                raise ValueError("update_chart.series must be a non-empty array")
            if len(categories) > 2_000 or len(series_items) > 100:
                raise ValueError("Chart data exceeds the supported size")
            data = CategoryChartData()
            data.categories = [str(value) for value in categories]
            for index, item in enumerate(series_items):
                if not isinstance(item, dict):
                    raise ValueError("Each update_chart series must be an object")
                values = item.get("values")
                if not isinstance(values, list) or len(values) != len(categories):
                    raise ValueError(
                        "Each chart series must contain one value for every category"
                    )
                normalized: list[float | None] = []
                for value in values:
                    if value is None:
                        normalized.append(None)
                        continue
                    try:
                        normalized.append(float(value))
                    except (TypeError, ValueError) as exc:
                        raise ValueError(f"Invalid chart value: {value}") from exc
                data.add_series(str(item.get("name") or f"Series {index + 1}"), normalized)
            try:
                chart.replace_data(data)
            except (AttributeError, ValueError) as exc:
                raise ValueError(
                    "This chart type does not support category-data replacement"
                ) from exc
            changed += 1
        if "show_legend" in operation:
            chart.has_legend = bool(operation["show_legend"])
            changed += 1
        colors = operation.get("series_colors")
        if colors is not None:
            if not isinstance(colors, list) or not colors:
                raise ValueError("series_colors must be a non-empty array")
            for index, series in enumerate(chart.series):
                color = colors[index % len(colors)]
                rgb = cls._color(color, f"series_colors[{index}]")
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = rgb
                except (AttributeError, ValueError):
                    pass
                try:
                    series.format.line.color.rgb = rgb
                except (AttributeError, ValueError):
                    pass
            changed += 1
        if changed == 0:
            raise ValueError("update_chart did not include any editable fields")
        return changed

    @classmethod
    def _apply_operations(
        cls,
        presentation: Any,
        operations: Iterable[Any],
        *,
        asset_paths: dict[str, Path] | None,
    ) -> int:
        changed = 0
        for operation in operations:
            if not isinstance(operation, dict):
                raise ValueError("Presentation operation 必须是对象")
            kind = str(operation.get("type") or "")
            if kind == "replace_text":
                old = str(operation.get("old") or "")
                if not old:
                    raise ValueError("replace_text.old 不能为空")
                changed += cls._replace_text(
                    presentation,
                    old,
                    str(operation.get("new") or ""),
                    replace_all=operation.get("all") is True,
                )
            elif kind == "replace_placeholders":
                values = operation.get("values")
                if not isinstance(values, dict) or not values:
                    raise ValueError("replace_placeholders.values 必须是非空对象")
                missing = []
                for key, value in values.items():
                    placeholder = str(key)
                    if not (placeholder.startswith("{{") and placeholder.endswith("}}")):
                        placeholder = "{{" + placeholder + "}}"
                    count = cls._replace_text(
                        presentation,
                        placeholder,
                        str(value),
                        replace_all=True,
                        required=False,
                    )
                    changed += count
                    if count == 0:
                        missing.append(placeholder)
                if missing and operation.get("allow_missing") is not True:
                    raise ValueError("演示文稿中没有找到占位符: " + ", ".join(missing))
            elif kind == "append_slides":
                slides = operation.get("slides")
                if not isinstance(slides, list) or not slides:
                    raise ValueError("append_slides.slides 必须是非空数组")
                if len(presentation.slides) + len(slides) > MAX_SLIDES:
                    raise ValueError(f"演示文稿不能超过 {MAX_SLIDES} 页")
                theme = cls._theme(operation.get("theme"))
                for slide_spec in slides:
                    cls._append_slide(
                        presentation,
                        slide_spec,
                        theme=theme,
                        asset_paths=asset_paths,
                    )
                    changed += 1
            elif kind == "update_slide":
                changed += cls._update_slide(presentation, operation, asset_paths)
            elif kind == "update_element":
                changed += cls._update_element(presentation, operation)
            elif kind == "update_table_cell":
                changed += cls._update_table_cell(presentation, operation)
            elif kind == "replace_image":
                changed += cls._replace_image(presentation, operation, asset_paths)
            elif kind == "update_chart":
                changed += cls._update_chart(presentation, operation)
            elif kind == "delete_element":
                cls._delete_element(presentation, operation)
                changed += 1
            elif kind == "delete_slide":
                index = cls._slide_index(presentation, operation.get("slide"))
                if len(presentation.slides) == 1:
                    raise ValueError("不能删除演示文稿的最后一页")
                cls._delete_slide(presentation, index)
                changed += 1
            elif kind == "move_slide":
                source = cls._slide_index(presentation, operation.get("slide"))
                target = cls._slide_index(presentation, operation.get("to"))
                cls._move_slide(presentation, source, target)
                changed += 1
            elif kind == "set_properties":
                cls._set_properties(presentation, operation)
                changed += 1
            else:
                raise ValueError(f"不支持的 Presentation operation: {kind}")
        return changed

    @staticmethod
    def _summary(presentation: Any) -> tuple[int, int, int, int, int]:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        text_count = 0
        picture_count = 0
        chart_count = 0
        note_count = 0
        characters = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    text_count += 1
                    characters += len(shape.text)
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1
                if getattr(shape, "has_chart", False):
                    chart_count += 1
            try:
                if not slide.has_notes_slide:
                    continue
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except (AttributeError, ValueError):
                notes = ""
            if notes:
                note_count += 1
                characters += len(notes)
        return text_count, picture_count, chart_count, note_count, characters
