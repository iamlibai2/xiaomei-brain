"""PPTX slide and speaker-note extraction."""

from __future__ import annotations

import re
from posixpath import dirname, join, normpath
from pathlib import Path
from zipfile import BadZipFile, ZipFile
from typing import Any

from xiaomei_brain.documents.models import DocumentExtraction, DocumentSection
from xiaomei_brain.documents.office_xml import bounded_text, read_xml


_DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_PRESENTATION_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_OFFICE_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_PACKAGE_REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"


def _member_number(member: str) -> int:
    match = re.search(r"(\d+)\.xml$", member)
    return int(match.group(1)) if match else 0


def _lines(root) -> list[str]:
    return [
        node.text.strip()
        for node in root.iter(f"{_DRAWING_NS}t")
        if node.text and node.text.strip()
    ]


def _relationships(archive: ZipFile, member: str) -> dict[str, tuple[str, str]]:
    if member not in archive.namelist():
        return {}
    root = read_xml(archive, member)
    return {
        str(node.get("Id")): (str(node.get("Target") or ""), str(node.get("Type") or ""))
        for node in root.iter(f"{_PACKAGE_REL_NS}Relationship")
        if node.get("Id")
    }


def _ordered_slides(archive: ZipFile, members: list[str]) -> list[str]:
    presentation_member = "ppt/presentation.xml"
    relationships_member = "ppt/_rels/presentation.xml.rels"
    if presentation_member in members and relationships_member in members:
        relations = _relationships(archive, relationships_member)
        root = read_xml(archive, presentation_member)
        ordered = []
        for node in root.iter(f"{_PRESENTATION_NS}sldId"):
            relation_id = node.get(f"{_OFFICE_REL_NS}id")
            target = relations.get(str(relation_id), ("", ""))[0]
            if not target:
                continue
            member = normpath(join(dirname(presentation_member), target))
            if member in members:
                ordered.append(member)
        if ordered:
            return ordered
    return sorted(
        (name for name in members if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
        key=_member_number,
    )


def _notes_member(archive: ZipFile, slide_member: str) -> str | None:
    slide_name = slide_member.rsplit("/", 1)[-1]
    relationships_member = (
        f"{dirname(slide_member)}/_rels/{slide_name}.rels"
    )
    for target, relation_type in _relationships(archive, relationships_member).values():
        if relation_type.endswith("/notesSlide"):
            return normpath(join(dirname(slide_member), target))
    return None


def _element_type(shape: Any) -> str:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return "line"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "has_chart", False):
            return "chart"
        if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}:
            return "shape"
        if getattr(shape, "has_text_frame", False):
            return "text"
    except Exception:
        pass
    return "object"


def _shape_summary(shape: Any) -> str:
    text = ""
    try:
        if getattr(shape, "has_text_frame", False):
            text = str(shape.text or "").strip()
        elif getattr(shape, "has_table", False):
            text = " / ".join(
                str(cell.text or "").strip()
                for row in shape.table.rows
                for cell in row.cells
                if str(cell.text or "").strip()
            )
        elif getattr(shape, "has_chart", False) and shape.chart.has_title:
            text = str(shape.chart.chart_title.text_frame.text or "").strip()
    except Exception:
        text = ""
    return re.sub(r"\s+", " ", text)[:240]


def _shape_index_lines(
    shapes: Any,
    *,
    slide_number: int,
    parent_ids: tuple[int, ...] = (),
    remaining: list[int],
) -> list[str]:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return []
    lines: list[str] = []
    for shape in shapes:
        if remaining[0] <= 0:
            break
        shape_ids = (*parent_ids, int(shape.shape_id))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            lines.extend(_shape_index_lines(
                shape.shapes,
                slide_number=slide_number,
                parent_ids=shape_ids,
                remaining=remaining,
            ))
            continue
        remaining[0] -= 1
        element_id = (
            f"slide-{slide_number}-shape-id-"
            + ".".join(str(value) for value in shape_ids)
        )
        position = tuple(
            round(float(getattr(shape, field, 0) or 0) / 360000, 2)
            for field in ("left", "top", "width", "height")
        )
        name = re.sub(r"\s+", " ", str(getattr(shape, "name", "") or "")).strip()[:120]
        summary = _shape_summary(shape)
        line = (
            f'- element_id="{element_id}" type={_element_type(shape)} '
            f'name="{name}" position_cm={position}'
        )
        if summary:
            line += f' text="{summary}"'
        lines.append(line)
    return lines


def _presentation_element_index(path: Path) -> dict[int, list[str]]:
    """Return a compact, best-effort index the Agent can use for precise edits."""
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
    except Exception:
        return {}
    result: dict[int, list[str]] = {}
    for slide_number, slide in enumerate(presentation.slides, start=1):
        result[slide_number] = _shape_index_lines(
            slide.shapes,
            slide_number=slide_number,
            remaining=[200],
        )
    return result


class PresentationExtractor:
    extractor_id = "document_presentation"
    extractor_version = "1.2.0"
    suffixes = (".pptx",)
    mime_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    def extract(self, path: Path) -> DocumentExtraction:
        element_index = _presentation_element_index(path)
        try:
            with ZipFile(path) as archive:
                members = archive.namelist()
                slides = _ordered_slides(archive, members)
                if not slides:
                    raise ValueError("演示文稿没有幻灯片")
                sections = []
                for index, member in enumerate(slides, start=1):
                    content = _lines(read_xml(archive, member))
                    note_member = _notes_member(archive, member)
                    if note_member is None:
                        # Keep best-effort support for minimal or damaged PPTX
                        # archives that omit slide relationship files.
                        note_member = (
                            "ppt/notesSlides/notesSlide"
                            f"{_member_number(member)}.xml"
                        )
                    if note_member and note_member in members:
                        note_lines = _lines(read_xml(archive, note_member))
                        if note_lines:
                            content.extend(["", "[备注]", *note_lines])
                    indexed = element_index.get(index, [])
                    if indexed:
                        content.extend(["", "[元素索引]", *indexed])
                    sections.append(DocumentSection(
                        key=f"slide:{index}",
                        title=f"第 {index} 页",
                        content=bounded_text("\n".join(content) or "[没有可提取的文字]"),
                        metadata={
                            "slide": index,
                            "element_count": len(indexed),
                            "element_index_truncated": len(indexed) >= 200,
                        },
                    ))
        except (BadZipFile, KeyError, ValueError) as exc:
            raise ValueError(f"无法解析演示文稿: {path.name}") from exc
        return DocumentExtraction(
            extractor_id=self.extractor_id,
            extractor_version=self.extractor_version,
            sections=tuple(sections),
            metadata={
                "format": "pptx",
                "slide_count": len(sections),
                "element_index": bool(element_index),
            },
        )
