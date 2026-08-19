"""Build and persist a browser-previewable presentation project.

The project deliberately follows PPTD's useful idea -- a small manifest plus
self-contained page descriptions -- without depending on Kimi's bundled web
editor.  JSON is valid YAML, so the generated ``.pptd`` and ``.page`` files
remain easy for agents and implementation staff to inspect and revise.
"""

from __future__ import annotations

import hashlib
import html
import json
import math
import mimetypes
import shutil
from pathlib import Path
from typing import Any, Iterable


PROJECT_SCHEMA = "xiaomei.presentation.v1"
PROJECT_GENERATOR_VERSION = 11
CANVAS_WIDTH = 960.0
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"


def presentation_project_directory(output_path: Path) -> Path:
    """Return the private companion-project directory for one PPTX output."""
    return output_path.parent / ".presentation" / output_path.stem


def build_presentation_project(pptx_path: Path, project_dir: Path) -> dict[str, Any]:
    """Convert a PPTX into a deterministic PPTD-style preview project."""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - writer already requires it
        raise ValueError("演示文稿预览依赖 python-pptx 未安装") from exc

    pptx_path = pptx_path.resolve(strict=True)
    project_dir = project_dir.resolve()
    if project_dir.parent.name != ".presentation":
        raise ValueError("演示文稿项目必须位于受控的 .presentation 目录")
    if project_dir.exists():
        shutil.rmtree(project_dir)
    pages_dir = project_dir / "pages"
    media_dir = project_dir / "media"
    pages_dir.mkdir(parents=True, exist_ok=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(pptx_path))
    slide_width = max(1, int(presentation.slide_width))
    slide_height = max(1, int(presentation.slide_height))
    canvas_width = CANVAS_WIDTH
    canvas_height = round(canvas_width * slide_height / slide_width, 3)
    scale_x = canvas_width / slide_width
    scale_y = canvas_height / slide_height
    title = str(presentation.core_properties.title or pptx_path.stem).strip() or pptx_path.stem
    source_revision = _source_revision(pptx_path)

    manifest_pages: list[str] = []
    preview_slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        page_path = f"pages/{slide_index:03d}.page"
        manifest_pages.append(page_path)
        elements: list[dict[str, Any]] = []
        image_counter = [0]
        background = _slide_background(
            slide,
            slide_index=slide_index,
            media_dir=media_dir,
            image_counter=image_counter,
        )
        for shape in slide.shapes:
            elements.extend(_shape_elements(
                shape,
                slide_index=slide_index,
                shape_id_path=(int(shape.shape_id),),
                scale_x=scale_x,
                scale_y=scale_y,
                media_dir=media_dir,
                image_counter=image_counter,
            ))
        page = {
            "pageType": "content",
            "background": background,
            "elements": elements,
        }
        _write_json(project_dir / page_path, page)
        preview_slides.append({
            "path": page_path,
            "index": slide_index,
            **page,
        })

    manifest_name = f"{pptx_path.stem}.pptd"
    manifest = {
        "version": "v2",
        "schema": PROJECT_SCHEMA,
        "title": title,
        "size": [canvas_width, canvas_height],
        "pages": manifest_pages,
        "sourceRevision": source_revision,
    }
    _write_json(project_dir / manifest_name, manifest)
    project = {
        "schema": PROJECT_SCHEMA,
        "generatorVersion": PROJECT_GENERATOR_VERSION,
        "title": title,
        "size": [canvas_width, canvas_height],
        "manifest": manifest_name,
        "sourceRevision": source_revision,
        "slides": preview_slides,
    }
    _write_json(project_dir / "project.json", project)
    return {
        "path": str(project_dir),
        "manifest": manifest_name,
        "slide_count": len(preview_slides),
        "schema": PROJECT_SCHEMA,
        "source_revision": source_revision,
    }


def _source_revision(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_presentation_project(project_dir: Path) -> dict[str, Any]:
    value = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    if value.get("schema") != PROJECT_SCHEMA or not isinstance(value.get("slides"), list):
        raise ValueError("演示文稿项目格式无效")
    media: dict[str, dict[str, str]] = {}
    for path in sorted((project_dir / "media").glob("*")):
        if not path.is_file():
            continue
        import base64
        relative = path.relative_to(project_dir).as_posix()
        media[relative] = {
            "mime_type": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
            "data_base64": base64.b64encode(path.read_bytes()).decode("ascii"),
        }
    value["media"] = media
    return value


def _shape_elements(
    shape: Any,
    *,
    slide_index: int,
    shape_id_path: tuple[int, ...],
    scale_x: float,
    scale_y: float,
    media_dir: Path,
    image_counter: list[int],
    coordinate_transform: tuple[float, float, float, float] = (1.0, 1.0, 0.0, 0.0),
    inherited_rotation: float = 0.0,
) -> list[dict[str, Any]]:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:  # pragma: no cover
        return []

    shape_type = shape.shape_type
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        nested: list[dict[str, Any]] = []
        child_transform = _group_child_transform(shape, coordinate_transform)
        group_rotation = inherited_rotation + _shape_rotation(shape)
        for child in shape.shapes:
            nested.extend(_shape_elements(
                child,
                slide_index=slide_index,
                shape_id_path=(*shape_id_path, int(child.shape_id)),
                scale_x=scale_x,
                scale_y=scale_y,
                media_dir=media_dir,
                image_counter=image_counter,
                coordinate_transform=child_transform,
                inherited_rotation=group_rotation,
            ))
        return nested

    transform_x, transform_y, offset_x, offset_y = coordinate_transform
    bounds = [
        round((offset_x + float(shape.left) * transform_x) * scale_x, 3),
        round((offset_y + float(shape.top) * transform_y) * scale_y, 3),
        round(float(shape.width) * transform_x * scale_x, 3),
        round(float(shape.height) * transform_y * scale_y, 3),
    ]
    element_id = (
        f"slide-{slide_index}-shape-id-"
        + ".".join(str(shape_id) for shape_id in shape_id_path)
    )

    formula = _formula_element(
        shape,
        element_id,
        bounds,
        rotation=inherited_rotation + _shape_rotation(shape),
    )
    if formula is not None:
        return [formula]

    if shape_type == MSO_SHAPE_TYPE.LINE:
        return [{
            "elementId": element_id,
            "elementType": "line",
            "bounds": bounds,
            "rotation": inherited_rotation + _shape_rotation(shape),
            "connectorKind": _connector_kind(shape),
            "flip": _shape_flip(shape),
            "adjustments": _shape_adjustments(shape),
            "line": _shape_line_with_theme(shape),
            "startArrow": _line_arrow(shape, "headEnd"),
            "endArrow": _line_arrow(shape, "tailEnd"),
        }]

    if (
        shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO}
        or shape._element.find(f".//{{{A_NS}}}videoFile") is not None
        or shape._element.find(f".//{{{A_NS}}}audioFile") is not None
        or shape._element.find(f".//{{{P14_NS}}}media") is not None
    ):
        media = _media_element(
            shape,
            element_id,
            bounds,
            rotation=inherited_rotation + _shape_rotation(shape),
            slide_index=slide_index,
            media_dir=media_dir,
            media_counter=image_counter,
        )
        return [media] if media is not None else []

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image_counter[0] += 1
            extension = str(shape.image.ext or "png").lower()
            filename = f"slide-{slide_index:03d}-{image_counter[0]:03d}.{extension}"
            (media_dir / filename).write_bytes(shape.image.blob)
            return [{
                "elementId": element_id,
                "elementType": "image",
                "bounds": bounds,
                "rotation": inherited_rotation + _shape_rotation(shape),
                "src": f"media/{filename}",
                "fit": {"mode": "fill"},
                "crop": _picture_crop(shape),
                "cropShape": _picture_crop_shape(shape),
                "shadow": _shape_shadow(shape),
            }]
        except Exception:
            return []

    if getattr(shape, "has_chart", False):
        return _chart_elements(
            shape,
            element_id,
            bounds,
            rotation=inherited_rotation + _shape_rotation(shape),
        )

    if getattr(shape, "has_table", False):
        return _table_elements(
            shape,
            element_id,
            bounds,
            rotation=inherited_rotation + _shape_rotation(shape),
        )

    text = ""
    if getattr(shape, "has_text_frame", False):
        text = str(shape.text or "").strip()
    style = _text_style(shape)
    rich_text = _rich_text(shape.text_frame) if getattr(shape, "has_text_frame", False) else None
    if shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}:
        fill = _shape_fill(
            shape,
            slide_index=slide_index,
            media_dir=media_dir,
            image_counter=image_counter,
        )
        line = _shape_line(shape)
        return [{
            "elementId": element_id,
            "elementType": "shape",
            "bounds": bounds,
            "rotation": inherited_rotation + _shape_rotation(shape),
            "shapeName": _shape_name(shape),
            "fill": fill,
            "line": line,
            "shadow": _shape_shadow(shape),
            "customGeometry": _custom_geometry(shape) if shape_type == MSO_SHAPE_TYPE.FREEFORM else None,
            "text": text,
            "textStyle": style,
            "richText": rich_text,
        }]
    if text:
        return [{
            "elementId": element_id,
            "elementType": "text",
            "bounds": bounds,
            "rotation": inherited_rotation + _shape_rotation(shape),
            "content": {"text": text, **style},
            "richText": rich_text,
            "shadow": _shape_shadow(shape),
        }]
    return []


def _formula_element(
    shape: Any,
    element_id: str,
    bounds: list[float],
    *,
    rotation: float = 0.0,
) -> dict[str, Any] | None:
    """Convert a native Office Math object into browser-native MathML."""
    try:
        equations = shape._element.findall(f".//{{{M_NS}}}oMath")
    except Exception:
        return None
    if not equations:
        return None
    rendered = [_omml_to_mathml(node) for node in equations]
    body = "<mspace linebreak=\"newline\"/>".join(value for value in rendered if value)
    if not body:
        return None
    fallback = " ".join(
        text.strip()
        for node in equations
        for text in ["".join(node.itertext())]
        if text.strip()
    )
    style = _text_style(shape)
    return {
        "elementId": element_id,
        "elementType": "formula",
        "bounds": bounds,
        "rotation": rotation,
        "mathMl": f'<math xmlns="http://www.w3.org/1998/Math/MathML" display="block"><mrow>{body}</mrow></math>',
        "fallbackText": fallback,
        "textStyle": style,
    }


def _omml_to_mathml(node: Any) -> str:
    """Render the common OMML structures produced by PowerPoint as MathML."""
    name = str(getattr(node, "tag", "")).rsplit("}", 1)[-1]
    children = list(node)

    def child(local_name: str) -> Any | None:
        return next(
            (item for item in children if str(item.tag).rsplit("}", 1)[-1] == local_name),
            None,
        )

    def render(local_name: str) -> str:
        item = child(local_name)
        return _omml_children(item) if item is not None else "<mrow/>"

    if name in {"oMath", "oMathPara", "e", "num", "den", "sub", "sup", "deg", "fName"}:
        return _omml_children(node)
    if name == "r":
        text = "".join(
            str(item.text or "") for item in node.findall(f".//{{{M_NS}}}t")
        )
        return _math_token(text)
    if name == "f":
        return f"<mfrac>{render('num')}{render('den')}</mfrac>"
    if name == "sSup":
        return f"<msup>{render('e')}{render('sup')}</msup>"
    if name == "sSub":
        return f"<msub>{render('e')}{render('sub')}</msub>"
    if name == "sSubSup":
        return f"<msubsup>{render('e')}{render('sub')}{render('sup')}</msubsup>"
    if name == "rad":
        degree = render("deg")
        base = render("e")
        return f"<mroot>{base}{degree}</mroot>" if _omml_has_content(child("deg")) else f"<msqrt>{base}</msqrt>"
    if name == "nary":
        properties = child("naryPr")
        operator = _omml_property(properties, "chr", "∫")
        base = f"<mo>{html.escape(operator)}</mo>"
        lower = render("sub")
        upper = render("sup")
        if _omml_has_content(child("sub")) or _omml_has_content(child("sup")):
            base = f"<munderover>{base}{lower}{upper}</munderover>"
        return f"<mrow>{base}{render('e')}</mrow>"
    if name == "d":
        properties = child("dPr")
        start = _omml_property(properties, "begChr", "(")
        end = _omml_property(properties, "endChr", ")")
        content = "".join(
            _omml_children(item)
            for item in children
            if str(item.tag).rsplit("}", 1)[-1] == "e"
        )
        return (
            f'<mfenced open="{html.escape(start, quote=True)}" '
            f'close="{html.escape(end, quote=True)}">{content}</mfenced>'
        )
    if name == "m":
        rows = []
        for row in children:
            if str(row.tag).rsplit("}", 1)[-1] != "mr":
                continue
            cells = "".join(
                f"<mtd>{_omml_children(item)}</mtd>"
                for item in row
                if str(item.tag).rsplit("}", 1)[-1] == "e"
            )
            rows.append(f"<mtr>{cells}</mtr>")
        return f"<mtable>{''.join(rows)}</mtable>"
    if name == "eqArr":
        rows = "".join(
            f"<mtr><mtd>{_omml_children(item)}</mtd></mtr>"
            for item in children
            if str(item.tag).rsplit("}", 1)[-1] == "e"
        )
        return f"<mtable>{rows}</mtable>"
    if name in {"limLow", "limUpp"}:
        tag = "munder" if name == "limLow" else "mover"
        limit_name = "lim"
        return f"<{tag}>{render('e')}{render(limit_name)}</{tag}>"
    if name == "acc":
        accent = _omml_property(child("accPr"), "chr", "ˆ")
        return f'<mover accent="true">{render("e")}<mo>{html.escape(accent)}</mo></mover>'
    if name == "bar":
        position = _omml_property(child("barPr"), "pos", "top")
        tag = "munder" if position == "bot" else "mover"
        return f"<{tag}>{render('e')}<mo>¯</mo></{tag}>"
    if name == "func":
        return f"<mrow>{render('fName')}<mo>⁡</mo>{render('e')}</mrow>"
    return _omml_children(node)


def _omml_children(node: Any | None) -> str:
    if node is None:
        return ""
    return "".join(
        _omml_to_mathml(item)
        for item in node
        if str(item.tag).rsplit("}", 1)[-1] not in {
            "ctrlPr", "rPr", "fPr", "radPr", "sSupPr", "sSubPr", "sSubSupPr",
            "naryPr", "dPr", "mPr", "eqArrPr", "limLowPr", "limUppPr", "accPr",
            "barPr", "funcPr",
        }
    )


def _math_token(value: str) -> str:
    value = str(value or "")
    if not value:
        return ""
    escaped = html.escape(value)
    stripped = value.strip()
    if stripped and all(character.isdigit() or character in ".," for character in stripped):
        return f"<mn>{escaped}</mn>"
    if len(stripped) == 1 and (stripped.isalpha() or "α" <= stripped <= "ω" or "Α" <= stripped <= "Ω"):
        return f"<mi>{escaped}</mi>"
    if stripped and all(character in "+−-×÷=*/≠≈<>≤≥±∓∝∞∂∆∇∈∉∪∩∧∨¬" for character in stripped):
        return f"<mo>{escaped}</mo>"
    return f"<mtext>{escaped}</mtext>"


def _omml_property(parent: Any | None, name: str, default: str) -> str:
    if parent is None:
        return default
    item = parent.find(f"{{{M_NS}}}{name}")
    if item is None:
        return default
    return str(item.get(f"{{{M_NS}}}val") or item.get("val") or default)


def _omml_has_content(node: Any | None) -> bool:
    return node is not None and bool("".join(node.itertext()).strip() or len(node))


def _media_element(
    shape: Any,
    element_id: str,
    bounds: list[float],
    *,
    rotation: float,
    slide_index: int,
    media_dir: Path,
    media_counter: list[int],
) -> dict[str, Any] | None:
    """Extract an embedded PowerPoint audio/video object and its poster frame."""
    root = shape._element
    video = root.find(f".//{{{A_NS}}}videoFile")
    audio = root.find(f".//{{{A_NS}}}audioFile")
    embedded = root.find(f".//{{{P14_NS}}}media")
    if video is None and audio is None and embedded is None:
        return None
    media_kind = "video" if video is not None else "audio" if audio is not None else "video"
    relationship_id = embedded.get(f"{{{R_NS}}}embed") if embedded is not None else None
    if not relationship_id:
        source_node = video if video is not None else audio
        relationship_id = source_node.get(f"{{{R_NS}}}link") if source_node is not None else None
    source = _write_related_media(
        shape.part,
        relationship_id,
        slide_index=slide_index,
        media_dir=media_dir,
        image_counter=media_counter,
    )
    poster_node = root.find(f".//{{{A_NS}}}blip")
    poster_id = poster_node.get(f"{{{R_NS}}}embed") if poster_node is not None else None
    poster = _write_related_media(
        shape.part,
        poster_id,
        slide_index=slide_index,
        media_dir=media_dir,
        image_counter=media_counter,
    )
    if not source and not poster:
        return None
    return {
        "elementId": element_id,
        "elementType": "media",
        "bounds": bounds,
        "rotation": rotation,
        "mediaKind": media_kind,
        "src": source,
        "posterSrc": poster,
        "mimeType": mimetypes.guess_type(source)[0] or (
            "video/mp4" if media_kind == "video" else "audio/mpeg"
        ),
    }


def _table_elements(
    shape: Any,
    element_id: str,
    bounds: list[float],
    *,
    rotation: float = 0.0,
) -> list[dict[str, Any]]:
    table = shape.table
    rows = len(table.rows)
    columns = len(table.columns)
    if rows <= 0 or columns <= 0:
        return []
    x, y, width, height = bounds
    cell_width = width / columns
    cell_height = height / rows
    values = []
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            values.append({
                "row": row_index,
                "column": column_index,
                "text": str(cell.text or ""),
                "fill": _fill_format(cell.fill),
                "textStyle": _text_frame_style(cell.text_frame),
                "columnSpan": int(getattr(cell, "span_width", 1) or 1),
                "rowSpan": int(getattr(cell, "span_height", 1) or 1),
                "hidden": bool(getattr(cell, "is_spanned", False)),
            })
    return [{
        "elementId": element_id,
        "elementType": "table",
        "bounds": bounds,
        "rotation": rotation,
        "rows": rows,
        "columns": columns,
        "cellWidth": cell_width,
        "cellHeight": cell_height,
        "cells": values,
    }]


def _chart_elements(
    shape: Any,
    element_id: str,
    bounds: list[float],
    *,
    rotation: float = 0.0,
) -> list[dict[str, Any]]:
    from pptx.oxml.ns import qn

    chart = shape.chart
    title = ""
    try:
        if chart.has_title:
            title = str(chart.chart_title.text_frame.text or "").strip()
    except Exception:
        pass
    categories: list[str] = []
    try:
        categories = [
            " / ".join(str(part) for part in label if part is not None)
            for label in chart.plots[0].categories.flattened_labels
        ]
    except Exception:
        pass
    chart_type = str(chart.chart_type or "chart").split("(", 1)[0].strip().lower()
    series_items: list[dict[str, Any]] = []
    for index, series in enumerate(chart.series):
        try:
            values = [_chart_value(value) for value in series.values]
        except Exception:
            values = []
        series_element = getattr(series, "_element", None)
        series_color = _chart_series_color(series, index)
        series_properties = series_element.find(qn("c:spPr")) if series_element is not None else None
        marker = _chart_marker(series_element)
        if "scatter" in chart_type and marker is None:
            marker = {"symbol": "circle", "size": 5.0}
        series_line = (
            _chart_inherited_line(
                series_properties,
                default_color=series_color,
                default_width=2.0,
            )
            if "line" in chart_type or "scatter" in chart_type or "radar" in chart_type
            else _xml_line(series_properties)
        )
        series_items.append({
            "name": str(series.name or f"Series {index + 1}"),
            "values": values,
            "xValues": _chart_cached_values(series_element, "xVal"),
            "color": series_color,
            "line": series_line,
            "marker": marker,
            "dataLabels": _chart_data_labels(series_element),
        })
    is_scatter = "scatter" in chart_type
    return [{
        "elementId": element_id,
        "elementType": "chart",
        "bounds": bounds,
        "rotation": rotation,
        "chartType": chart_type,
        "title": title,
        "titleStyle": _chart_title_style(chart),
        "fill": _chart_fill(chart),
        "categories": categories,
        "series": series_items,
        "hasLegend": bool(chart.has_legend),
        "legend": _chart_legend(chart),
        "plotArea": _chart_plot_area(chart),
        "gapWidth": _chart_number(chart, "gapWidth"),
        "overlap": _chart_number(chart, "overlap"),
        "roundedCorners": _chart_flag(chart._chartSpace.find(qn("c:roundedCorners"))),
        "categoryAxis": _chart_axis(chart, "valAx", positions={"b", "t"}) if is_scatter else _chart_axis(chart, "catAx"),
        "valueAxis": _chart_axis(chart, "valAx", positions={"l", "r"}) if is_scatter else _chart_axis(chart, "valAx"),
    }]


def _chart_cached_values(series_element: Any, value_name: str) -> list[float | str | None]:
    """Read cached chart values such as scatter-series X coordinates."""
    if series_element is None:
        return []
    try:
        from pptx.oxml.ns import qn

        container = series_element.find(qn(f"c:{value_name}"))
        if container is None:
            return []
        points = sorted(
            container.findall(f".//{qn('c:pt')}"),
            key=lambda point: int(point.get("idx") or 0),
        )
        return [
            _chart_value(point.find(qn("c:v")).text if point.find(qn("c:v")) is not None else None)
            for point in points
        ]
    except Exception:
        return []


def _chart_marker(series_element: Any) -> dict[str, Any] | None:
    if series_element is None:
        return None
    try:
        from pptx.oxml.ns import qn

        marker = series_element.find(qn("c:marker"))
        if marker is None:
            return None
        symbol = _element_value(marker.find(qn("c:symbol"))) or "none"
        if symbol == "none":
            return {"symbol": "none", "size": 0}
        size = _element_value(marker.find(qn("c:size")))
        properties = marker.find(qn("c:spPr"))
        fill = _xml_solid_fill(properties)
        return {
            "symbol": symbol,
            "size": float(size) if size else 5.0,
            "fill": fill,
            "line": _xml_line(properties),
        }
    except Exception:
        return None


def _chart_data_labels(series_element: Any) -> dict[str, Any] | None:
    if series_element is None:
        return None
    try:
        from pptx.oxml.ns import qn

        labels = series_element.find(qn("c:dLbls"))
        if labels is None and series_element.getparent() is not None:
            labels = series_element.getparent().find(qn("c:dLbls"))
        if labels is None:
            return None
        return {
            "showValue": _chart_flag(labels.find(qn("c:showVal"))),
            "showCategory": _chart_flag(labels.find(qn("c:showCatName"))),
            "showSeries": _chart_flag(labels.find(qn("c:showSerName"))),
            "showPercent": _chart_flag(labels.find(qn("c:showPercent"))),
            "position": _element_value(labels.find(qn("c:dLblPos"))) or "outEnd",
            "numberFormat": _element_value(labels.find(qn("c:numFmt")), "formatCode"),
            "style": _chart_text_properties(labels.find(qn("c:txPr"))),
        }
    except Exception:
        return None


def _chart_legend(chart: Any) -> dict[str, Any] | None:
    if not chart.has_legend:
        return None
    try:
        from pptx.oxml.ns import qn

        legend = chart._chartSpace.find(f".//{qn('c:legend')}")
        if legend is None:
            return None
        return {
            "position": _element_value(legend.find(qn("c:legendPos"))) or "r",
            "overlay": _chart_flag(legend.find(qn("c:overlay"))),
            "style": _chart_text_properties(legend.find(qn("c:txPr"))),
            "layout": _chart_manual_layout(legend.find(qn("c:layout"))),
        }
    except Exception:
        return None


def _chart_plot_area(chart: Any) -> dict[str, Any] | None:
    try:
        from pptx.oxml.ns import qn

        plot_area = chart._chartSpace.find(f".//{qn('c:plotArea')}")
        return _chart_manual_layout(plot_area.find(qn("c:layout"))) if plot_area is not None else None
    except Exception:
        return None


def _chart_manual_layout(layout: Any) -> dict[str, Any] | None:
    if layout is None:
        return None
    try:
        from pptx.oxml.ns import qn

        manual = layout.find(qn("c:manualLayout"))
        if manual is None:
            return None
        result: dict[str, Any] = {
            "target": _element_value(manual.find(qn("c:layoutTarget"))) or "outer",
        }
        for name in ("x", "y", "w", "h"):
            value = _element_value(manual.find(qn(f"c:{name}")))
            if value is not None:
                result[name] = float(value)
        return result if len(result) > 1 else None
    except Exception:
        return None


def _chart_number(chart: Any, element_name: str) -> float | None:
    try:
        from pptx.oxml.ns import qn

        value = _element_value(chart._chartSpace.find(f".//{qn(f'c:{element_name}')}"))
        return float(value) if value is not None else None
    except Exception:
        return None


def _chart_flag(element: Any) -> bool:
    if element is None:
        return False
    return str(element.get("val", "1")).lower() not in {"0", "false", "off"}


def _xml_solid_fill(properties: Any) -> dict[str, Any]:
    if properties is None:
        return {"type": "none", "color": "transparent"}
    solid = properties.find(f"{{{A_NS}}}solidFill")
    if solid is None:
        return {"type": "none", "color": "transparent"}
    return {"type": "solid", "color": _xml_color(solid) or "transparent"}


def _chart_fill(chart: Any) -> dict[str, Any]:
    """Return the chart-area fill from ``c:chartSpace/c:spPr``."""
    try:
        from pptx.oxml.ns import qn

        properties = chart._chartSpace.find(qn("c:spPr"))
        if properties is None or properties.find(qn("a:noFill")) is not None:
            return {"type": "none", "color": "transparent"}
        solid = properties.find(qn("a:solidFill"))
        if solid is not None:
            color = solid.find(qn("a:srgbClr"))
            if color is not None and color.get("val"):
                return {"type": "solid", "color": f"#{color.get('val')}"}
    except Exception:
        pass
    return {"type": "unknown", "color": "transparent"}


def _chart_title_style(chart: Any) -> dict[str, Any]:
    try:
        if chart.has_title:
            return _text_frame_style(chart.chart_title.text_frame)
    except Exception:
        pass
    return {"fontSize": 18, "color": "#253047", "bold": False, "align": "center"}


def _chart_axis(
    chart: Any,
    axis_name: str,
    *,
    positions: set[str] | None = None,
) -> dict[str, Any]:
    """Extract native axis visibility, labels, line and major gridline style."""
    try:
        from pptx.oxml.ns import qn

        axes = chart._chartSpace.findall(f".//{qn(f'c:{axis_name}')}")
        axis = next(
            (
                candidate
                for candidate in axes
                if positions is None
                or _element_value(candidate.find(qn("c:axPos"))) in positions
            ),
            None,
        )
        if axis is None:
            return {"visible": False}
        deleted = axis.find(qn("c:delete"))
        label_position = axis.find(qn("c:tickLblPos"))
        visible = not (deleted is not None and deleted.get("val") == "1")
        scaling = axis.find(qn("c:scaling"))
        result: dict[str, Any] = {
            "visible": visible,
            "labelsVisible": not (
                label_position is not None and label_position.get("val") == "none"
            ),
            "position": _element_value(axis.find(qn("c:axPos"))),
            "line": _chart_inherited_line(
                axis.find(qn("c:spPr")),
                default_color="#888888",
                default_width=0.75,
            ),
            "labelStyle": _chart_text_properties(axis.find(qn("c:txPr"))),
            "numberFormat": _element_value(axis.find(qn("c:numFmt")), "formatCode"),
        }
        if scaling is not None:
            for source_name, target_name in (("min", "minimum"), ("max", "maximum")):
                value = _element_value(scaling.find(qn(f"c:{source_name}")))
                if value not in (None, ""):
                    try:
                        result[target_name] = float(value)
                    except ValueError:
                        pass
        major_unit = _element_value(axis.find(qn("c:majorUnit")))
        if major_unit not in (None, ""):
            try:
                result["majorUnit"] = float(major_unit)
            except ValueError:
                pass
        gridlines = axis.find(qn("c:majorGridlines"))
        result["majorGridline"] = (
            _chart_inherited_line(
                gridlines.find(qn("c:spPr")),
                default_color="#D9DEE7",
                default_width=0.5,
            )
            if gridlines is not None
            else {"type": "none", "color": "transparent", "width": 0}
        )
        return result
    except Exception:
        return {"visible": False}


def _chart_inherited_line(
    properties: Any,
    *,
    default_color: str,
    default_width: float,
) -> dict[str, Any]:
    """Resolve chart lines while preserving Office theme inheritance.

    Missing ``c:spPr`` / ``a:ln`` means "use the chart theme default", not
    "hide the line".  An explicit ``a:noFill`` still means no line.
    """
    if properties is None:
        return {"type": "solid", "color": default_color, "width": default_width}
    try:
        line = properties.find(f"{{{A_NS}}}ln")
        if line is None:
            return {"type": "solid", "color": default_color, "width": default_width}
        if line.find(f"{{{A_NS}}}noFill") is not None:
            return {"type": "none", "color": "transparent", "width": 0}
        parsed = _xml_line(properties)
        if parsed.get("type") in {"solid", "dash", "dot"} and parsed.get("color") != "transparent":
            return parsed
        width = round(float(line.get("w") or default_width * 12700) / 12700, 2)
        dash = line.find(f"{{{A_NS}}}prstDash")
        dash_value = str(dash.get("val")) if dash is not None and dash.get("val") else "solid"
        line_type = "dash" if "dash" in dash_value else "dot" if "dot" in dash_value else "solid"
        return {"type": line_type, "color": default_color, "width": width}
    except Exception:
        return {"type": "solid", "color": default_color, "width": default_width}


def _element_value(element: Any, attribute: str = "val") -> str | None:
    if element is None:
        return None
    value = element.get(attribute)
    return str(value) if value is not None else None


def _chart_text_properties(text_properties: Any) -> dict[str, Any]:
    style: dict[str, Any] = {
        "fontSize": 10,
        "color": "#697386",
        "bold": False,
        "fontFamily": "Arial",
    }
    if text_properties is None:
        return style
    try:
        run_properties = text_properties.find(f".//{{{A_NS}}}defRPr")
        if run_properties is None:
            run_properties = text_properties.find(f".//{{{A_NS}}}rPr")
        if run_properties is None:
            return style
        size = run_properties.get("sz")
        if size:
            style["fontSize"] = round(float(size) / 100, 2)
        if run_properties.get("b") is not None:
            style["bold"] = run_properties.get("b") in {"1", "true"}
        color = _xml_color(run_properties)
        if color:
            style["color"] = color
        for tag in ("ea", "latin", "cs"):
            font = run_properties.find(f"{{{A_NS}}}{tag}")
            if font is not None and font.get("typeface"):
                style["fontFamily"] = str(font.get("typeface"))
                break
    except Exception:
        pass
    return style


def _xml_line(properties: Any) -> dict[str, Any]:
    if properties is None:
        return {"type": "none", "color": "transparent", "width": 0}
    try:
        line = properties.find(f"{{{A_NS}}}ln")
        if line is None or line.find(f"{{{A_NS}}}noFill") is not None:
            return {"type": "none", "color": "transparent", "width": 0}
        explicit_fill = next(
            (
                line.find(f"{{{A_NS}}}{fill_name}")
                for fill_name in ("solidFill", "gradFill", "pattFill", "blipFill")
                if line.find(f"{{{A_NS}}}{fill_name}") is not None
            ),
            None,
        )
        if explicit_fill is None:
            return {"type": "none", "color": "transparent", "width": 0}
        color = _xml_color(explicit_fill)
        if not color:
            return {"type": "unknown", "color": "transparent", "width": 0}
        width = round(float(line.get("w") or 12700) / 12700, 2)
        dash = line.find(f"{{{A_NS}}}prstDash")
        dash_value = str(dash.get("val")) if dash is not None and dash.get("val") else "solid"
        line_type = "dash" if "dash" in dash_value else "dot" if "dot" in dash_value else "solid"
        return {"type": line_type, "color": color, "width": width}
    except Exception:
        return {"type": "none", "color": "transparent", "width": 0}


def _chart_value(value: Any) -> float | str | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    try:
        return float(value)
    except (TypeError, ValueError):
        return str(value)


def _chart_series_color(series: Any, index: int) -> str:
    palette = ("#4F6BED", "#16A085", "#F39C12", "#E15B64", "#7A5AF8", "#3498DB")
    try:
        from pptx.oxml.ns import qn

        series_element = getattr(series, "_element", None)
        properties = series_element.find(qn("c:spPr")) if series_element is not None else None
        if properties is not None:
            fill_color = _xml_color(properties.find(qn("a:solidFill")))
            if fill_color:
                return fill_color
            line = properties.find(qn("a:ln"))
            line_color = _xml_color(line)
            if line_color:
                return line_color
    except Exception:
        pass
    return palette[index % len(palette)]


def _text_style(shape: Any) -> dict[str, Any]:
    try:
        return _text_frame_style(shape.text_frame)
    except Exception:
        return {"fontSize": 18, "color": "#253047", "bold": False, "align": "left"}


def _text_frame_style(text_frame: Any) -> dict[str, Any]:
    style: dict[str, Any] = {
        "fontSize": 18,
        "color": "#253047",
        "bold": False,
        "align": "left",
        "verticalAlign": _vertical_alignment(text_frame),
        "margins": _text_frame_margins(text_frame),
    }
    try:
        paragraphs: Iterable[Any] = text_frame.paragraphs
        paragraph = next(iter(paragraphs), None)
        if paragraph is None:
            return style
        alignment = str(paragraph.alignment or "").lower()
        if "center" in alignment:
            style["align"] = "center"
        elif "right" in alignment:
            style["align"] = "right"
        run = next(iter(paragraph.runs), None)
        font = run.font if run is not None else paragraph.font
        if font.size is not None:
            style["fontSize"] = round(float(font.size.pt), 2)
        family = _font_family(font)
        if family:
            style["fontFamily"] = family
        if font.bold is not None:
            style["bold"] = bool(font.bold)
        color = _font_color(font)
        if color:
            style["color"] = color
    except Exception:
        pass
    return style


def _rich_text(text_frame: Any) -> dict[str, Any]:
    """Preserve paragraph and run formatting instead of flattening a text box."""
    paragraphs: list[dict[str, Any]] = []
    try:
        for paragraph in text_frame.paragraphs:
            paragraph_item: dict[str, Any] = {
                "align": _paragraph_alignment(paragraph),
                "level": int(paragraph.level or 0),
                "bullet": _paragraph_has_bullet(paragraph),
                "runs": [],
            }
            for field, value in (
                ("spaceBefore", _length_points(paragraph.space_before)),
                ("spaceAfter", _length_points(paragraph.space_after)),
            ):
                if value is not None:
                    paragraph_item[field] = value
            line_height = _paragraph_line_height(paragraph)
            if line_height is not None:
                paragraph_item["lineHeight"] = line_height

            runs: list[dict[str, Any]] = []
            for run in paragraph.runs:
                if not run.text:
                    continue
                item: dict[str, Any] = {"text": str(run.text)}
                font = run.font
                if font.size is not None:
                    item["fontSize"] = round(float(font.size.pt), 2)
                family = _font_family(font)
                if family:
                    item["fontFamily"] = family
                if font.bold is not None:
                    item["bold"] = bool(font.bold)
                if font.italic is not None:
                    item["italic"] = bool(font.italic)
                color = _font_color(font)
                if color:
                    item["color"] = color
                runs.append(item)
            if not runs and paragraph.text:
                runs.append({"text": str(paragraph.text)})
            paragraph_item["runs"] = runs
            paragraphs.append(paragraph_item)
    except Exception:
        return {"paragraphs": []}
    return {"paragraphs": paragraphs}


def _paragraph_alignment(paragraph: Any) -> str:
    alignment = str(getattr(paragraph, "alignment", "") or "").lower()
    if "center" in alignment:
        return "center"
    if "right" in alignment:
        return "right"
    if "justify" in alignment or "distributed" in alignment:
        return "justify"
    return "left"


def _vertical_alignment(text_frame: Any) -> str:
    anchor = str(getattr(text_frame, "vertical_anchor", "") or "").lower()
    if "middle" in anchor:
        return "middle"
    if "bottom" in anchor:
        return "bottom"
    return "top"


def _text_frame_margins(text_frame: Any) -> list[float]:
    values = []
    for name in ("margin_top", "margin_right", "margin_bottom", "margin_left"):
        value = _length_points(getattr(text_frame, name, None))
        values.append(value if value is not None else 0.0)
    return values


def _length_points(value: Any) -> float | None:
    try:
        if value is None:
            return None
        return round(float(value.pt), 2)
    except Exception:
        return None


def _paragraph_line_height(paragraph: Any) -> dict[str, float] | None:
    try:
        value = paragraph.line_spacing
        if value is None:
            return None
        if hasattr(value, "pt"):
            return {"points": round(float(value.pt), 2)}
        return {"multiple": round(float(value), 3)}
    except Exception:
        return None


def _paragraph_has_bullet(paragraph: Any) -> bool:
    try:
        properties = paragraph._p.pPr
        if properties is None:
            return False
        return any(
            child.tag.rsplit("}", 1)[-1] in {"buChar", "buAutoNum", "buBlip"}
            for child in properties
        )
    except Exception:
        return False


def _font_family(font: Any) -> str:
    try:
        properties = font._rPr
        if properties is not None:
            for tag in ("ea", "latin", "cs"):
                node = properties.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
                if node is not None and node.get("typeface"):
                    return str(node.get("typeface"))
    except Exception:
        pass
    try:
        return str(font.name) if font.name else ""
    except Exception:
        return ""


def _shape_rotation(shape: Any) -> float:
    try:
        return round(float(shape.rotation or 0), 3)
    except Exception:
        return 0.0


def _font_color(font: Any) -> str:
    try:
        rgb = font.color.rgb
        return f"#{rgb}" if rgb is not None else ""
    except Exception:
        return ""


def _shape_fill(
    shape: Any,
    *,
    slide_index: int,
    media_dir: Path,
    image_counter: list[int],
) -> dict[str, Any]:
    """Return the real fill without inventing a white background."""
    try:
        parsed = _xml_fill(
            shape._element.spPr,
            part=shape.part,
            slide_index=slide_index,
            media_dir=media_dir,
            image_counter=image_counter,
        )
        if parsed is not None:
            return parsed
    except Exception:
        pass
    return _fill_format(shape.fill)


def _fill_format(fill: Any) -> dict[str, Any]:
    try:
        fill_type = fill.type
        if fill_type is None or "background" in str(fill_type).lower():
            return {"type": "none", "color": "transparent"}
        rgb = fill.fore_color.rgb
        if rgb is not None:
            return {"type": "solid", "color": f"#{rgb}"}
    except Exception:
        pass
    return {"type": "unknown", "color": "transparent"}


def _shape_line(shape: Any) -> dict[str, Any]:
    """Return only an explicitly visible line; never synthesize a border."""
    try:
        properties = shape._element.spPr
        if properties.find(f"{{{A_NS}}}ln") is not None:
            return _xml_line(properties)
    except Exception:
        pass
    try:
        fill_type = shape.line.fill.type
        if fill_type is None or "background" in str(fill_type).lower():
            return {"type": "none", "color": "transparent", "width": 0}
        rgb = shape.line.color.rgb
        if rgb is not None:
            line_width = shape.line.width
            width = (
                max(0.0, round(float(line_width.pt), 2))
                if line_width is not None
                else 1.0
            )
            return {"type": "solid", "color": f"#{rgb}", "width": width}
    except Exception:
        pass
    return {"type": "unknown", "color": "transparent", "width": 0}


def _shape_line_with_theme(shape: Any) -> dict[str, Any]:
    """Resolve a connector's inherited theme line without affecting ordinary shapes."""
    explicit = _shape_line(shape)
    if explicit.get("width", 0) > 0 and explicit.get("color") != "transparent":
        return explicit
    try:
        style = shape._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
        )
        line_reference = style.find(f"{{{A_NS}}}lnRef") if style is not None else None
        if line_reference is None:
            return explicit
        color_node = next(iter(line_reference), None)
        scheme_name = (
            str(color_node.get("val") or "accent1")
            if color_node is not None
            else "accent1"
        )
        color = _theme_scheme_color(shape.part, scheme_name) or "#4472C4"
        theme_line = _theme_line_style(shape.part, int(line_reference.get("idx") or 1))
        return {
            "type": theme_line.get("type", "solid"),
            "color": color,
            "width": theme_line.get("width", 1.0),
        }
    except Exception:
        return explicit


def _theme_root(part: Any) -> Any | None:
    try:
        from lxml import etree

        master_part = part.slide_layout.slide_master.part
        theme_part = next(
            rel.target_part
            for rel in master_part.rels.values()
            if str(rel.reltype).endswith("/theme")
        )
        return etree.fromstring(theme_part.blob)
    except Exception:
        return None


def _theme_scheme_color(part: Any, scheme_name: str) -> str:
    root = _theme_root(part)
    if root is None:
        return ""
    aliases = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}
    name = aliases.get(scheme_name, scheme_name)
    node = root.find(f".//{{{A_NS}}}clrScheme/{{{A_NS}}}{name}")
    return _xml_color(node) if node is not None else ""


def _theme_line_style(part: Any, index: int) -> dict[str, Any]:
    root = _theme_root(part)
    if root is None:
        return {"type": "solid", "width": 1.0}
    lines = root.findall(
        f".//{{{A_NS}}}fmtScheme/{{{A_NS}}}lnStyleLst/{{{A_NS}}}ln"
    )
    if not lines:
        return {"type": "solid", "width": 1.0}
    line = lines[max(0, min(len(lines) - 1, index - 1))]
    width = round(float(line.get("w") or 12700) / 12700, 2)
    dash = line.find(f"{{{A_NS}}}prstDash")
    dash_value = str(dash.get("val") or "solid") if dash is not None else "solid"
    return {
        "type": "dash" if "dash" in dash_value else "dot" if "dot" in dash_value else "solid",
        "width": width,
    }


def _connector_kind(shape: Any) -> str:
    try:
        geometry = shape._element.spPr.find(f"{{{A_NS}}}prstGeom")
        return str(geometry.get("prst") or "line") if geometry is not None else "line"
    except Exception:
        return "line"


def _shape_flip(shape: Any) -> list[bool]:
    try:
        transform = shape._element.spPr.find(f"{{{A_NS}}}xfrm")
        return [
            str(transform.get("flipH") or "0").lower() in {"1", "true"},
            str(transform.get("flipV") or "0").lower() in {"1", "true"},
        ] if transform is not None else [False, False]
    except Exception:
        return [False, False]


def _shape_adjustments(shape: Any) -> list[float]:
    try:
        geometry = shape._element.spPr.find(f"{{{A_NS}}}prstGeom")
        values: list[float] = []
        guides = geometry.findall(f".//{{{A_NS}}}gd") if geometry is not None else []
        for guide in guides:
            formula = str(guide.get("fmla") or "")
            if formula.startswith("val "):
                values.append(float(formula[4:]) / 100000)
        return values
    except Exception:
        return []


def _line_arrow(shape: Any, tag_name: str) -> dict[str, Any] | None:
    try:
        line = shape._element.spPr.find(f"{{{A_NS}}}ln")
        arrow = line.find(f"{{{A_NS}}}{tag_name}") if line is not None else None
        arrow_type = str(arrow.get("type") or "none") if arrow is not None else "none"
        if arrow_type == "none":
            return None
        return {
            "type": arrow_type,
            "width": str(arrow.get("w") or "med"),
            "length": str(arrow.get("len") or "med"),
        }
    except Exception:
        return None


def _custom_geometry(shape: Any) -> dict[str, Any] | None:
    try:
        geometry = shape._element.spPr.find(f"{{{A_NS}}}custGeom")
        if geometry is None:
            return None
        rendered_paths: list[dict[str, Any]] = []
        paths = geometry.findall(f".//{{{A_NS}}}pathLst/{{{A_NS}}}path")
        for path in paths:
            width = max(1.0, float(path.get("w") or shape.width or 1))
            height = max(1.0, float(path.get("h") or shape.height or 1))
            commands: list[str] = []
            for command in path:
                command_name = str(command.tag).rsplit("}", 1)[-1]
                points = command.findall(f"{{{A_NS}}}pt")
                coordinates = [
                    (float(point.get("x") or 0), float(point.get("y") or 0))
                    for point in points
                ]
                if command_name == "moveTo" and coordinates:
                    commands.append(f"M {coordinates[0][0]} {coordinates[0][1]}")
                elif command_name == "lnTo" and coordinates:
                    commands.append(f"L {coordinates[0][0]} {coordinates[0][1]}")
                elif command_name == "cubicBezTo" and len(coordinates) >= 3:
                    commands.append("C " + " ".join(f"{x} {y}" for x, y in coordinates[:3]))
                elif command_name == "quadBezTo" and len(coordinates) >= 2:
                    commands.append("Q " + " ".join(f"{x} {y}" for x, y in coordinates[:2]))
                elif command_name == "close":
                    commands.append("Z")
            if commands:
                rendered_paths.append({
                    "d": " ".join(commands),
                    "viewBox": [width, height],
                    "fill": str(path.get("fill") or "norm") != "none",
                    "stroke": str(path.get("stroke") or "1").lower() not in {"0", "false"},
                })
        return {"paths": rendered_paths} if rendered_paths else None
    except Exception:
        return None


def _slide_background(
    slide: Any,
    *,
    slide_index: int,
    media_dir: Path,
    image_counter: list[int],
) -> dict[str, Any]:
    try:
        from pptx.oxml.ns import qn

        background = slide._element.cSld.find(qn("p:bg"))
        properties = background.find(qn("p:bgPr")) if background is not None else None
        parsed = _xml_fill(
            properties,
            part=slide.part,
            slide_index=slide_index,
            media_dir=media_dir,
            image_counter=image_counter,
        )
        if parsed is not None:
            return parsed
    except Exception:
        pass
    try:
        rgb = slide.background.fill.fore_color.rgb
        return {"type": "solid", "color": f"#{rgb}" if rgb is not None else "#FFFFFF"}
    except Exception:
        return {"type": "solid", "color": "#FFFFFF"}


def _xml_fill(
    properties: Any,
    *,
    part: Any,
    slide_index: int,
    media_dir: Path,
    image_counter: list[int],
) -> dict[str, Any] | None:
    if properties is None:
        return None
    if properties.find(f"{{{A_NS}}}noFill") is not None:
        return {"type": "none", "color": "transparent"}
    solid = properties.find(f"{{{A_NS}}}solidFill")
    if solid is not None:
        return {"type": "solid", "color": _xml_color(solid) or "transparent"}
    gradient = properties.find(f"{{{A_NS}}}gradFill")
    if gradient is not None:
        stops = []
        for stop in gradient.findall(f".//{{{A_NS}}}gs"):
            color = _xml_color(stop)
            if not color:
                continue
            stops.append({
                "position": round(float(stop.get("pos") or 0) / 100000, 4),
                "color": color,
            })
        if len(stops) >= 2:
            linear = gradient.find(f"{{{A_NS}}}lin")
            angle = round(float(linear.get("ang") or 0) / 60000, 3) if linear is not None else 0.0
            return {
                "type": "gradient",
                "gradientType": "linear" if linear is not None else "radial",
                "angle": angle,
                "stops": stops,
            }
    image_fill = properties.find(f"{{{A_NS}}}blipFill")
    if image_fill is not None:
        blip = image_fill.find(f"{{{A_NS}}}blip")
        relationship_id = blip.get(f"{{{R_NS}}}embed") if blip is not None else None
        source = _write_related_media(
            part,
            relationship_id,
            slide_index=slide_index,
            media_dir=media_dir,
            image_counter=image_counter,
        )
        if source:
            return {
                "type": "image",
                "src": source,
                "fit": {"mode": "cover"},
                "crop": _xml_crop(image_fill),
            }
    return None


def _write_related_media(
    part: Any,
    relationship_id: str | None,
    *,
    slide_index: int,
    media_dir: Path,
    image_counter: list[int],
) -> str:
    if not relationship_id:
        return ""
    try:
        related = part.related_part(relationship_id)
        extension = Path(str(related.partname)).suffix or mimetypes.guess_extension(related.content_type) or ".png"
        image_counter[0] += 1
        filename = f"slide-{slide_index:03d}-{image_counter[0]:03d}{extension.lower()}"
        (media_dir / filename).write_bytes(related.blob)
        return f"media/{filename}"
    except Exception:
        return ""


def _xml_crop(image_fill: Any) -> dict[str, float]:
    source_rect = image_fill.find(f"{{{A_NS}}}srcRect") if image_fill is not None else None
    if source_rect is None:
        return {"left": 0.0, "top": 0.0, "right": 0.0, "bottom": 0.0}
    return {
        name: round(float(source_rect.get(attribute) or 0) / 100000, 6)
        for name, attribute in (("left", "l"), ("top", "t"), ("right", "r"), ("bottom", "b"))
    }


def _picture_crop(shape: Any) -> dict[str, float]:
    return {
        name: round(float(getattr(shape, f"crop_{name}", 0) or 0), 6)
        for name in ("left", "top", "right", "bottom")
    }


def _picture_crop_shape(shape: Any) -> str:
    try:
        geometry = shape._element.spPr.find(f"{{{A_NS}}}prstGeom")
        value = str(geometry.get("prst") or "rect") if geometry is not None else "rect"
        if value in {"ellipse", "oval"}:
            return "ellipse"
        if "round" in value.lower():
            return "roundRect"
    except Exception:
        pass
    return "rect"


def _shape_shadow(shape: Any) -> dict[str, Any] | None:
    try:
        shadow = shape._element.spPr.find(f".//{{{A_NS}}}outerShdw")
        if shadow is None:
            return None
        blur = round(float(shadow.get("blurRad") or 0) / 12700, 2)
        distance = float(shadow.get("dist") or 0) / 12700
        direction = math.radians(float(shadow.get("dir") or 0) / 60000)
        return {
            "blur": blur,
            "color": _xml_color(shadow) or "#00000040",
            "offset": [
                round(distance * math.cos(direction), 2),
                round(distance * math.sin(direction), 2),
            ],
        }
    except Exception:
        return None


def _xml_color(parent: Any) -> str:
    if parent is None:
        return ""
    for tag in ("srgbClr", "sysClr", "schemeClr", "prstClr"):
        node = parent.find(f".//{{{A_NS}}}{tag}")
        if node is None:
            continue
        value = node.get("val") or node.get("lastClr")
        if tag == "schemeClr" and not node.get("lastClr"):
            value = {"dk1": "000000", "lt1": "FFFFFF"}.get(str(value), "")
        if not value or len(str(value)) != 6:
            continue
        alpha = node.find(f"{{{A_NS}}}alpha")
        if alpha is not None and alpha.get("val"):
            opacity = max(0, min(255, round(float(alpha.get("val")) / 100000 * 255)))
            return f"#{str(value).upper()}{opacity:02X}"
        return f"#{str(value).upper()}"
    return ""


def _group_child_transform(
    shape: Any,
    parent: tuple[float, float, float, float],
) -> tuple[float, float, float, float]:
    parent_x_scale, parent_y_scale, parent_x_offset, parent_y_offset = parent
    try:
        xfrm = shape._element.grpSpPr.xfrm
        child_offset_x = float(xfrm.chOff.x)
        child_offset_y = float(xfrm.chOff.y)
        child_extent_x = max(1.0, float(xfrm.chExt.cx))
        child_extent_y = max(1.0, float(xfrm.chExt.cy))
        group_x_scale = float(shape.width) / child_extent_x
        group_y_scale = float(shape.height) / child_extent_y
        return (
            parent_x_scale * group_x_scale,
            parent_y_scale * group_y_scale,
            parent_x_offset + parent_x_scale * (float(shape.left) - child_offset_x * group_x_scale),
            parent_y_offset + parent_y_scale * (float(shape.top) - child_offset_y * group_y_scale),
        )
    except Exception:
        return parent


def _shape_name(shape: Any) -> str:
    try:
        name = str(shape.auto_shape_type or "").lower()
    except Exception:
        return "rect"
    if "ellipse" in name or "oval" in name:
        return "ellipse"
    if "round" in name:
        return "roundRect"
    return "rect"


def _write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(value, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
