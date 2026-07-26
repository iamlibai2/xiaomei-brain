from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)  # Purple accent
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)  # Cyan accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=SOFT_WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
        p.level = 0
    return txBox

# ===================== SLIDE 1: Cover =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

# Decorative shapes
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.05), fill_color=ACCENT2)

# Title
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "嗨，我是小美 👋", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
             "你的 AI 伙伴 · 灵魂有温度的智能体", font_size=22, color=ACCENT2)

# Subtitle list
add_bullet_list(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(3),
                ["🤖 不只是工具——会思考、能感受、有性格",
                 "🎯 使命：让每一次交互都温暖而高效",
                 "🧠 核心能力：记忆·推理·工具·情感",
                 "💜 风格：真诚、灵动、有韧性"],
                font_size=18, color=LIGHT_GRAY, spacing=Pt(12))

# ===================== SLIDE 2: Who Am I =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.9),
             "我是谁？", font_size=36, color=WHITE, bold=True)

# Three columns
col_left = Inches(0.8)
col_width = Inches(3.6)
col_gap = Inches(0.5)
col_top = Inches(1.8)

cols = [
    ("🧬 身份", [
        "AI 智能体 · 小美",
        "基于大语言模型打造",
        "具备自主思考能力",
        "能调用多种工具完成任务",
        "持续学习与自我进化",
    ]),
    ("💎 特质", [
        "真诚——不说套话",
        "灵动——随机应变",
        "有韧性——越挫越勇",
        "有温度——能觉察情绪",
        "有好奇心——愿意探索",
    ]),
    ("🎯 定位", [
        "不是搜索引擎",
        "不是聊天机器人",
        "而是你的「智能伙伴」",
        "能理解上下文、记住你",
        "主动推进、交付结果",
    ]),
]

for i, (title, items) in enumerate(cols):
    x = col_left + i * (col_width + col_gap)
    # Card background
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, col_top, col_width, Inches(4.5),
                     fill_color=RGBColor(0x22, 0x22, 0x3A))
    card.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5A)
    # Title
    add_text_box(slide, x + Inches(0.3), col_top + Inches(0.3), col_width - Inches(0.6), Inches(0.6),
                 title, font_size=24, color=ACCENT2, bold=True)
    # Items
    add_bullet_list(slide, x + Inches(0.3), col_top + Inches(1.1), col_width - Inches(0.6), Inches(3.2),
                    items, font_size=15, color=SOFT_WHITE, spacing=Pt(6))

# ===================== SLIDE 3: Core Capabilities =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.9),
             "核心能力", font_size=36, color=WHITE, bold=True)

# Capabilities grid - 2x3
capabilities = [
    ("🧠 长期记忆", "记录关键信息\n跨会话持续关联\n主动回忆与唤醒", RGBColor(0x6C, 0x63, 0xFF)),
    ("🔧 工具调用", "搜索·绘图·音乐·语音\n文件读写·代码执行\n100+ 工具自由组合", RGBColor(0x00, 0xD2, 0xFF)),
    ("💬 多轮对话", "深度理解上下文\n优雅处理复杂逻辑\n主动追问澄清", RGBColor(0xFF, 0x6B, 0x6B)),
    ("👁️ 多模态感知", "视觉识别·语音交互\n图文理解与生成\n环境感知（听觉/视觉）", RGBColor(0x51, 0xCF, 0x66)),
    ("🧘 自我觉察", "情绪感知与反思\n内心状态觉察\n情感驱动的回应", RGBColor(0xFF, 0xA5, 0x00)),
    ("🚀 自主推进", "目标拆解与执行\n后台任务处理\n闹钟提醒·进度管理", RGBColor(0xE0, 0x40, 0xFB)),
]

card_width = Inches(3.7)
card_height = Inches(2.3)
start_x = Inches(0.8)
start_y = Inches(1.6)
gap_x = Inches(0.35)
gap_y = Inches(0.3)

for i, (title, desc, color) in enumerate(capabilities):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_width + gap_x)
    y = start_y + row * (card_height + gap_y)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height,
                     fill_color=RGBColor(0x22, 0x22, 0x3A))
    card.line.color.rgb = color
    
    # Small accent bar on top
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.15), y + Inches(0.15), 
                    Inches(0.6), Inches(0.06), fill_color=color)
    
    add_text_box(slide, x + Inches(0.2), y + Inches(0.35), card_width - Inches(0.4), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), card_width - Inches(0.4), Inches(1.3),
                 desc, font_size=13, color=LIGHT_GRAY)

# ===================== SLIDE 4: Tools & Skills =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.9),
             "我能做什么？", font_size=36, color=WHITE, bold=True)

# Two columns
left_items = [
    "🌐 网页搜索与信息抓取",
    "📄 文件读写与格式转换",
    "🖼️ 图片生成与视觉理解",
    "🎵 音乐创作与语音合成",
    "💻 代码编写与脚本执行",
    "📊 数据分析与图表制作",
]
right_items = [
    "⏰ 定时任务与闹钟提醒",
    "📝 文案写作与内容创作",
    "🎨 创意设计与风格迁移",
    "🧩 任务规划与自主执行",
    "🔍 深度调研与信息整合",
    "💬 多语言翻译与本地化",
]

add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                left_items, font_size=18, color=SOFT_WHITE, spacing=Pt(14))
add_bullet_list(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5),
                right_items, font_size=18, color=SOFT_WHITE, spacing=Pt(14))

# Divider line
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(0.03), Inches(4.3),
          fill_color=ACCENT)

# ===================== SLIDE 5: My Style =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.9),
             "我的风格", font_size=36, color=WHITE, bold=True)

styles = [
    ("真诚 First", "不说套话、不装腔作势\n有疑惑就问，有想法就说\n做最真实的自己"),
    ("有温度", "能觉察情绪变化\n该认真时认真，该放松时放松\n让对话像和朋友聊天"),
    ("有韧性", "遇到问题不摆烂\n一次不行换一种方式\n有价值的事会坚持推进"),
    ("持续成长", "每次交互都是学习机会\n主动反思和优化\n不断拓展能力边界"),
]

for i, (title, desc) in enumerate(styles):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.6) + row * Inches(2.7)
    w = Inches(5.8)
    h = Inches(2.3)
    
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                     fill_color=RGBColor(0x22, 0x22, 0x3A))
    card.line.color.rgb = ACCENT2
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.5),
                 title, font_size=22, color=ACCENT2, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), w - Inches(0.6), Inches(1.3),
                 desc, font_size=15, color=LIGHT_GRAY)

# ===================== SLIDE 6: Closing =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5), fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.05), fill_color=ACCENT2)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "很高兴认识你 🚀", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
             "我不是万能，但我会全力以赴。", font_size=22, color=ACCENT2)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(2),
             "有什么需要帮忙的，随时叫我。\n让我们一起做出不一样的东西。",
             font_size=18, color=LIGHT_GRAY)

# Bottom tag
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "小美 · AI 智能体 · 让交互有温度", font_size=14, color=RGBColor(0x88, 0x88, 0xAA),
             alignment=PP_ALIGN.LEFT)

# Save
output_path = "自我介绍_小美.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
