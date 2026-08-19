import base64
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

from xiaomei_brain.plugin.context import PluginContext
from xiaomei_brain.plugin.registry import PluginRegistry
from xiaomei_brain.plugins.tools.document_presentation.adapter import (
    register as register_presentation,
)
from xiaomei_brain.plugins.tools.document_presentation.extractor import (
    PresentationExtractor,
)
from xiaomei_brain.plugins.tools.document_io.tool import create_write_document_tool
from xiaomei_brain.tools.execution_context import bind_tool_execution


PNG_1PX = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def _presentation_registry() -> PluginRegistry:
    registry = PluginRegistry()
    context = PluginContext({}, "document_presentation", "test", registry)
    register_presentation(context)
    return registry


def test_presentation_plugin_owns_writer_and_skill_directory():
    registry = _presentation_registry()

    assert registry.get_document_writer("presentation") is not None
    assert registry.list_document_writers() == ["presentation"]
    skill_dirs = registry.get_skill_directories()
    assert len(skill_dirs) == 1
    assert (Path(skill_dirs[0]) / "SKILL.md").is_file()


def test_write_document_creates_themed_presentation_with_image_and_notes(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    assets = workspace / "work"
    assets.mkdir(parents=True)
    (assets / "cover.png").write_bytes(PNG_1PX)
    spec = workspace / "deck.json"
    spec.write_text(json.dumps({
        "properties": {"author": "Xiaomei", "title": "Product"},
        "page": {"size": "wide"},
        "theme": {
            "background_color": "F7F9FC",
            "title_color": "172033",
            "text_color": "354052",
            "accent_color": "4F6BED",
            "font_family": "Microsoft YaHei",
        },
        "slides": [
            {
                "type": "title",
                "title": "产品介绍",
                "subtitle": "让工作自然流动",
                "notes": "开场说明",
            },
            {
                "type": "content",
                "title": "核心能力",
                "bullets": [
                    "理解真实意图",
                    {"text": "跨渠道连续关系", "level": 1},
                ],
            },
            {
                "type": "image",
                "title": "产品界面",
                "image": {
                    "workspace_path": "work/cover.png",
                    "x_cm": 6,
                    "y_cm": 4,
                    "width_cm": 20,
                },
            },
        ],
    }, ensure_ascii=False), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-create",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="deck.json",
            output_name="product.pptx",
        )

    assert result.get("success") is True, result
    assert result["validation"]["valid"] is True
    assert result["validation"]["delivery_ready"] is True
    assert isinstance(result["validation"]["issues"], list)
    assert result["validation"]["slide_count"] == 3
    assert result["validation"]["picture_count"] == 1
    assert result["validation"]["note_slide_count"] == 1
    assert result["presentation_project"]["schema"] == "xiaomei.presentation.v1"
    assert result["presentation_project"]["slide_count"] == 3
    project_dir = outputs / ".presentation" / "product"
    assert (project_dir / "product.pptd").is_file()
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    assert len(project["slides"]) == 3
    assert any(
        element["elementType"] == "image"
        for element in project["slides"][2]["elements"]
    )
    deck = Presentation(outputs / "product.pptx")
    assert round(deck.slide_width / Cm(1), 2) == 33.87
    assert round(deck.slide_height / Cm(1), 2) == 19.05
    assert deck.core_properties.author == "Xiaomei"
    assert deck.slides[0].notes_slide.notes_text_frame.text == "开场说明"
    assert "• 理解真实意图" in "\n".join(
        shape.text for shape in deck.slides[1].shapes if shape.has_text_frame
    )


def test_presentation_validation_returns_structured_quality_issues(tmp_path):
    from xiaomei_brain.documents.presentation_project import build_presentation_project
    from xiaomei_brain.plugins.tools.document_presentation.validator import (
        validate_presentation_project,
    )

    source = tmp_path / "quality-issues.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    small = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(6), Cm(1))
    small.text = "难以阅读的小字"
    small.text_frame.paragraphs[0].runs[0].font.size = Cm(0.2)
    outside = slide.shapes.add_textbox(Cm(30), Cm(3), Cm(8), Cm(2))
    outside.text = "超出页面"
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "quality-issues"
    build_presentation_project(source, project_dir)
    validation = validate_presentation_project(project_dir)

    assert validation["valid"] is False
    assert validation["delivery_ready"] is False
    assert validation["error_count"] >= 1
    assert validation["warning_count"] >= 1
    codes = {issue["code"] for issue in validation["issues"]}
    assert "out_of_bounds" in codes
    assert "font_too_small" in codes
    for issue in validation["issues"]:
        assert set(("page", "element_id", "severity", "reason", "suggestion")) <= issue.keys()


def test_write_document_creates_and_updates_native_scatter_and_radar_charts(tmp_path):
    from xiaomei_brain.plugins.tools.document_presentation.writer import PresentationWriter

    source = tmp_path / "scatter-radar-created.pptx"
    specification = {
        "slides": [
            {
                "type": "blank",
                "elements": [{
                    "type": "chart",
                    "chart_type": "scatter",
                    "x_cm": 2,
                    "y_cm": 2,
                    "width_cm": 20,
                    "height_cm": 11,
                    "title": "投入与产出",
                    "series": [{
                        "name": "项目",
                        "x_values": [1, 2, 4],
                        "values": [4, 9, 16],
                    }],
                }],
            },
            {
                "type": "blank",
                "elements": [{
                    "type": "chart",
                    "chart_type": "radar_filled",
                    "x_cm": 2,
                    "y_cm": 2,
                    "width_cm": 20,
                    "height_cm": 11,
                    "title": "能力雷达",
                    "categories": ["质量", "速度", "成本", "服务"],
                    "series": [{"name": "当前", "values": [82, 68, 74, 91]}],
                }],
            },
        ],
    }

    result = PresentationWriter().write(specification, source)

    assert result["validation"]["chart_count"] == 2
    deck = Presentation(source)
    scatter = deck.slides[0].shapes[0].chart
    radar = deck.slides[1].shapes[0].chart
    assert scatter.chart_type == XL_CHART_TYPE.XY_SCATTER
    assert radar.chart_type == XL_CHART_TYPE.RADAR_FILLED
    project = json.loads(
        (tmp_path / ".presentation" / "scatter-radar-created" / "project.json")
        .read_text(encoding="utf-8")
    )
    scatter_element = project["slides"][0]["elements"][0]
    radar_element = project["slides"][1]["elements"][0]
    assert scatter_element["series"][0]["xValues"] == [1.0, 2.0, 4.0]
    assert scatter_element["series"][0]["values"] == [4.0, 9.0, 16.0]
    assert radar_element["categories"] == ["质量", "速度", "成本", "服务"]

    updated = tmp_path / "scatter-updated.pptx"
    scatter_id = scatter_element["elementId"]
    radar_id = radar_element["elementId"]
    update_result = PresentationWriter().write(
        {
            "operations": [
                {
                    "type": "update_chart",
                    "slide": 1,
                    "element_id": scatter_id,
                    "series": [{
                        "name": "修订项目",
                        "x_values": [2, 3, 5],
                        "values": [5, 10, 20],
                    }],
                },
                {
                    "type": "update_chart",
                    "slide": 2,
                    "element_id": radar_id,
                    "categories": ["质量", "速度", "成本", "服务"],
                    "series": [{"name": "目标", "values": [90, 80, 70, 95]}],
                },
            ],
        },
        updated,
        source_path=source,
    )

    assert updated.is_file()
    assert updated.stat().st_size > 0
    assert updated != source
    assert update_result["validation"]["changed_items"] == 2
    updated_project = json.loads(
        (tmp_path / ".presentation" / "scatter-updated" / "project.json")
        .read_text(encoding="utf-8")
    )
    updated_scatter = updated_project["slides"][0]["elements"][0]
    updated_radar = updated_project["slides"][1]["elements"][0]
    assert updated_scatter["series"][0]["xValues"] == [2.0, 3.0, 5.0]
    assert updated_scatter["series"][0]["values"] == [5.0, 10.0, 20.0]
    assert updated_radar["series"][0]["name"] == "目标"
    assert updated_radar["series"][0]["values"] == [90.0, 80.0, 70.0, 95.0]


def test_presentation_validation_rejects_chart_data_length_mismatch(tmp_path):
    from xiaomei_brain.plugins.tools.document_presentation.validator import (
        validate_presentation_project,
    )

    project_dir = tmp_path / ".presentation" / "bad-chart"
    project_dir.mkdir(parents=True)
    (project_dir / "project.json").write_text(json.dumps({
        "schema": "xiaomei.presentation.v1",
        "size": [960, 540],
        "slides": [{
            "index": 2,
            "background": {"type": "solid", "color": "#FFFFFF"},
            "elements": [{
                "elementId": "slide-2-shape-id-7",
                "elementType": "chart",
                "bounds": [80, 80, 600, 320],
                "categories": ["Q1", "Q2", "Q3"],
                "series": [{"name": "Sales", "values": [10, 20]}],
            }],
        }],
    }), encoding="utf-8")

    validation = validate_presentation_project(project_dir)

    issue = next(
        item for item in validation["issues"]
        if item["code"] == "chart_data_length_mismatch"
    )
    assert issue["page"] == 2
    assert issue["element_id"] == "slide-2-shape-id-7"
    assert issue["severity"] == "error"


def test_write_document_creates_native_shape_line_table_and_chart(tmp_path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    spec = workspace / "native-elements.json"
    spec.write_text(json.dumps({
        "page": {"size": "wide"},
        "theme": {
            "accent_color": "2F6B4F",
            "text_color": "172033",
            "font_family": "Microsoft YaHei",
        },
        "slides": [{
            "type": "blank",
            "elements": [
                {
                    "type": "shape",
                    "shape": "round_rect",
                    "name": "ApprovalNode",
                    "x_cm": 1,
                    "y_cm": 1,
                    "width_cm": 7,
                    "height_cm": 2.5,
                    "fill_color": "EAF2EC",
                    "line_color": "2F6B4F",
                    "line_width_pt": 2,
                    "text": "方案确认",
                    "font_size_pt": 18,
                    "bold": True,
                    "align": "center",
                    "vertical": "middle",
                },
                {
                    "type": "line",
                    "connector": "elbow",
                    "name": "ApprovalFlow",
                    "x_cm": 8,
                    "y_cm": 2.25,
                    "to_x_cm": 12,
                    "to_y_cm": 4,
                    "line_color": "2F6B4F",
                    "line_width_pt": 2.5,
                    "line_dash": "dash",
                    "end_arrow": {"type": "triangle", "width": "lg", "length": "lg"},
                },
                {
                    "type": "table",
                    "name": "SalesTable",
                    "x_cm": 1,
                    "y_cm": 5,
                    "width_cm": 13,
                    "height_cm": 6,
                    "column_widths_cm": [5, 4, 4],
                    "header_style": {
                        "fill_color": "2F6B4F",
                        "text_color": "FFFFFF",
                        "bold": True,
                        "align": "center",
                    },
                    "cell_style": {"font_size_pt": 13},
                    "data": [
                        ["地区", "一季度", "二季度"],
                        ["华东", 120, 148],
                        ["华南", 98, {"text": 126, "bold": True}],
                    ],
                },
                {
                    "type": "chart",
                    "chart_type": "column",
                    "name": "SalesChart",
                    "x_cm": 15,
                    "y_cm": 1,
                    "width_cm": 17,
                    "height_cm": 10,
                    "title": "季度销售额",
                    "categories": ["一季度", "二季度"],
                    "series": [
                        {"name": "华东", "values": [120, 148]},
                        {"name": "华南", "values": [98, 126]},
                    ],
                    "show_legend": True,
                    "legend_position": "bottom",
                    "series_colors": ["2F6B4F", "C6F24E"],
                    "show_values": True,
                },
            ],
        }],
    }, ensure_ascii=False), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-native-elements",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="native-elements.json",
            output_name="native-elements.pptx",
        )

    assert result.get("success") is True, result
    assert result["validation"]["chart_count"] == 1
    deck = Presentation(outputs / "native-elements.pptx")
    shapes = list(deck.slides[0].shapes)
    node = next(shape for shape in shapes if shape.name == "ApprovalNode")
    assert node.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert node.text == "方案确认"
    connector = next(shape for shape in shapes if shape.name == "ApprovalFlow")
    assert connector.shape_type == MSO_SHAPE_TYPE.LINE
    line = connector._element.spPr.find(qn("a:ln"))
    assert line.find(qn("a:prstDash")).get("val") == "dash"
    assert line.find(qn("a:tailEnd")).get("type") == "triangle"
    table_shape = next(shape for shape in shapes if shape.name == "SalesTable")
    assert table_shape.has_table
    assert table_shape.table.cell(1, 0).text == "华东"
    assert table_shape.table.cell(2, 2).text == "126"
    chart_shape = next(shape for shape in shapes if shape.name == "SalesChart")
    assert chart_shape.has_chart
    assert chart_shape.chart.chart_title.text_frame.text == "季度销售额"
    assert list(chart_shape.chart.series[0].values) == [120.0, 148.0]
    project = json.loads(
        (outputs / ".presentation" / "native-elements" / "project.json").read_text(
            encoding="utf-8"
        )
    )
    element_types = {item["elementType"] for item in project["slides"][0]["elements"]}
    assert {"shape", "line", "table", "chart"}.issubset(element_types)
    extracted = PresentationExtractor().extract(outputs / "native-elements.pptx")
    content = extracted.sections[0].content
    for element_type in ("shape", "line", "table", "chart"):
        assert f"type={element_type}" in content


def test_blank_slide_does_not_render_top_level_title(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    spec = workspace / "blank-title.json"
    title = "01 - Agent designed canvas"
    spec.write_text(json.dumps({
        "slides": [{
            "type": "blank",
            "title": title,
            "elements": [{
                "type": "text",
                "name": "CanvasTitle",
                "x_cm": 1.2,
                "y_cm": 0.5,
                "width_cm": 30,
                "height_cm": 0.8,
                "text": title,
                "font_size_pt": 22,
                "bold": True,
            }],
        }],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-blank-title",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="blank-title.json",
            output_name="blank-title.pptx",
        )

    assert result.get("success") is True, result
    slide = Presentation(outputs / "blank-title.pptx").slides[0]
    assert not any(shape.name == "XiaomeiTitle" for shape in slide.shapes)
    assert [
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text == title
    ] == [title]


def test_write_document_revises_presentation_copy_and_preserves_source(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "template.pptx"
    original = Presentation()
    slide = original.slides.add_slide(original.slide_layouts[1])
    slide.shapes.title.text = "原方案"
    slide.placeholders[1].text = "客户：{{customer_name}}"
    original.save(source)
    source_bytes = source.read_bytes()
    spec = workspace / "revise.json"
    spec.write_text(json.dumps({
        "operations": [
            {
                "type": "replace_placeholders",
                "values": {"customer_name": "星海科技"},
            },
            {
                "type": "update_slide",
                "slide": 1,
                "title": "更新方案",
                "notes": "重点说明客户价值",
            },
            {
                "type": "append_slides",
                "slides": [
                    {"type": "section", "title": "实施计划"},
                    {"type": "content", "title": "下一步", "bullets": ["启动试点"]},
                ],
            },
            {"type": "move_slide", "slide": 3, "to": 2},
            {"type": "delete_slide", "slide": 3},
            {"type": "set_properties", "author": "Xiaomei", "title": "最终方案"},
        ],
    }, ensure_ascii=False), encoding="utf-8")
    attachment = {
        "id": "deck-source",
        "name": "template.pptx",
        "kind": "document",
        "local_path": str(source),
    }

    with bind_tool_execution(
        tool_call_id="call-presentation-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=(attachment,),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="revise.json",
            output_name="updated.pptx",
            source_attachment_id="deck-source",
        )

    assert result.get("success") is True, result
    assert result["validation"]["slide_count"] == 2
    assert source.read_bytes() == source_bytes
    updated = Presentation(outputs / "updated.pptx")
    assert updated.core_properties.author == "Xiaomei"
    assert updated.core_properties.title == "最终方案"
    assert updated.slides[0].shapes.title.text == "更新方案"
    assert "星海科技" in "\n".join(
        shape.text for shape in updated.slides[0].shapes if shape.has_text_frame
    )
    assert updated.slides[0].notes_slide.notes_text_frame.text == "重点说明客户价值"
    assert any(
        shape.text == "下一步"
        for shape in updated.slides[1].shapes
        if shape.has_text_frame
    )
    extracted = PresentationExtractor().extract(outputs / "updated.pptx")
    assert "更新方案" in extracted.sections[0].content
    assert "下一步" in extracted.sections[1].content
    assert "重点说明客户价值" in extracted.sections[0].content


def test_presentation_writer_rejects_missing_image_and_invalid_slide(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    missing_image = workspace / "missing-image.json"
    missing_image.write_text(json.dumps({
        "slides": [{
            "type": "image",
            "title": "图片",
            "image": {"workspace_path": "work/missing.png"},
        }],
    }, ensure_ascii=False), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-invalid",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="missing-image.json",
            output_name="invalid.pptx",
        )

    assert "error" in result
    assert not list(outputs.glob("*.pptx"))


def test_write_document_updates_exact_presentation_element_and_rebuilds_project(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "source.pptx"
    original = Presentation()
    slide = original.slides.add_slide(original.slide_layouts[6])
    slide.shapes.add_textbox(Cm(2), Cm(2), Cm(12), Cm(3)).text = "Original text"
    original.save(source)
    spec = workspace / "element-update.json"
    spec.write_text(json.dumps({
        "operations": [{
            "type": "update_element",
            "slide": 1,
            "element_id": "slide-1-shape-1",
            "text": "Updated text",
            "text_color": "336699",
            "fill_color": "F5F7FA",
            "font_size_pt": 24,
            "bold": True,
        }],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-element-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=({
            "id": "source-deck",
            "name": "source.pptx",
            "kind": "document",
            "local_path": str(source),
        },),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="element-update.json",
            output_name="edited.pptx",
            source_attachment_id="source-deck",
        )

    assert result.get("success") is True, result
    edited = Presentation(outputs / "edited.pptx")
    shape = edited.slides[0].shapes[0]
    assert shape.text == "Updated text"
    assert shape.text_frame.paragraphs[0].runs[0].font.bold is True
    project_dir = outputs / ".presentation" / "edited"
    assert result["presentation_project"]["path"] == str(project_dir)
    assert (project_dir / "edited.pptd").is_file()
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    assert project["slides"][0]["elements"][0]["content"]["text"] == "Updated text"
    assert project["slides"][0]["elements"][0]["elementId"].startswith(
        "slide-1-shape-id-"
    )
    assert len(project["sourceRevision"]) == 64
    assert sorted(path.name for path in (outputs / ".presentation").iterdir()) == ["edited"]


def test_presentation_annotation_operations_update_table_cell_and_replace_image(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "source.pptx"
    replacement = tmp_path / "replacement.png"
    replacement.write_bytes(PNG_1PX)

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Cm(1), Cm(1), Cm(12), Cm(4))
    table_shape.table.cell(1, 1).text = "Old value"
    image_shape = slide.shapes.add_picture(str(replacement), Cm(15), Cm(1), width=Cm(3))
    deck.save(source)

    from xiaomei_brain.documents.presentation_project import build_presentation_project

    preview_dir = tmp_path / ".presentation" / "source"
    build_presentation_project(source, preview_dir)
    project = json.loads((preview_dir / "project.json").read_text(encoding="utf-8"))
    elements = project["slides"][0]["elements"]
    table_id = next(item["elementId"] for item in elements if item["elementType"] == "table")
    image_id = next(item["elementId"] for item in elements if item["elementType"] == "image")
    assert table_id.endswith(str(table_shape.shape_id))
    assert image_id.endswith(str(image_shape.shape_id))

    spec = workspace / "precise-update.json"
    spec.write_text(json.dumps({
        "operations": [
            {
                "type": "update_table_cell",
                "slide": 1,
                "element_id": table_id,
                "row": 2,
                "column": 2,
                "text": "New value",
                "fill_color": "EAF2EC",
                "bold": True,
            },
            {
                "type": "replace_image",
                "slide": 1,
                "element_id": image_id,
                "attachment_id": "replacement-image",
            },
        ],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-precise-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=(
            {"id": "source-deck", "name": "source.pptx", "kind": "document", "local_path": str(source)},
            {"id": "replacement-image", "name": "replacement.png", "kind": "image", "local_path": str(replacement)},
        ),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="precise-update.json",
            output_name="precise.pptx",
            source_attachment_id="source-deck",
        )

    assert result.get("success") is True, result
    updated = Presentation(outputs / "precise.pptx")
    assert updated.slides[0].shapes[0].table.cell(1, 1).text == "New value"
    assert updated.slides[0].shapes[0].table.cell(1, 1).text_frame.paragraphs[0].runs[0].font.bold is True


def test_write_document_rejects_stale_presentation_annotation(tmp_path):
    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "source.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(2)).text = "Hello"
    deck.save(source)
    spec = workspace / "update.json"
    spec.write_text(json.dumps({
        "operations": [{"type": "update_element", "slide": 1, "element_id": "slide-1-shape-1", "text": "Changed"}],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-stale",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=({
            "id": "source-deck",
            "name": "source.pptx",
            "kind": "document",
            "local_path": str(source),
            "annotation": {
                "kind": "presentation",
                "source_revision": "0" * 64,
            },
        },),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="update.json",
            output_name="stale.pptx",
            source_attachment_id="source-deck",
        )

    assert result["subtype"] == "stale_presentation_selection"
    assert not (outputs / "stale.pptx").exists()


def test_presentation_chart_is_previewable_and_remains_native_when_updated(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Pt

    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "chart-source.pptx"

    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("East", [10, 18, 24])
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(2), Cm(2), Cm(22), Cm(12),
        data,
    )
    chart_shape.chart.has_title = True
    chart_shape.chart.chart_title.text_frame.text = "Quarterly sales"
    chart_shape.chart.has_legend = True
    chart_shape.chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart_shape.chart.legend.font.color.rgb = RGBColor(0x61, 0x68, 0x78)
    chart_shape.chart.legend.font.size = Pt(9)
    plot = chart_shape.chart.plots[0]
    plot.gap_width = 175
    plot.overlap = 20
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    plot.data_labels.font.size = Pt(10)
    chart_shape.chart.category_axis.format.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    chart_shape.chart.category_axis.format.line.width = Pt(1)
    chart_shape.chart.category_axis.tick_labels.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    chart_shape.chart.category_axis.tick_labels.font.size = Pt(10)
    chart_shape.chart.value_axis.format.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    chart_shape.chart.value_axis.format.line.width = Pt(1)
    chart_shape.chart.value_axis.tick_labels.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    chart_shape.chart.value_axis.tick_labels.font.size = Pt(10)
    chart_shape.chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x3A, 0x3A, 0x3A)
    chart_shape.chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    rounded = OxmlElement("c:roundedCorners")
    rounded.set("val", "1")
    chart_shape.chart._chartSpace.append(rounded)
    layout = OxmlElement("c:layout")
    manual = OxmlElement("c:manualLayout")
    for name, value in (("x", "0.12"), ("y", "0.10"), ("w", "0.72"), ("h", "0.68")):
        node = OxmlElement(f"c:{name}")
        node.set("val", value)
        manual.append(node)
    layout.append(manual)
    chart_shape.chart._chartSpace.find(f".//{qn('c:plotArea')}").insert(0, layout)
    deck.save(source)

    from xiaomei_brain.documents.presentation_project import build_presentation_project

    preview_dir = tmp_path / ".presentation" / "chart-source"
    build_presentation_project(source, preview_dir)
    project = json.loads((preview_dir / "project.json").read_text(encoding="utf-8"))
    chart_element = next(
        item for item in project["slides"][0]["elements"]
        if item["elementType"] == "chart"
    )
    assert chart_element["elementId"].endswith(str(chart_shape.shape_id))
    assert chart_element["title"] == "Quarterly sales"
    assert chart_element["fill"]["color"] == "transparent"
    assert chart_element["categories"] == ["Q1", "Q2", "Q3"]
    assert chart_element["series"][0]["values"] == [10.0, 18.0, 24.0]
    assert chart_element["categoryAxis"]["visible"] is True
    assert chart_element["categoryAxis"]["labelsVisible"] is True
    assert chart_element["categoryAxis"]["line"] == {
        "type": "solid", "color": "#888888", "width": 1.0,
    }
    assert chart_element["categoryAxis"]["labelStyle"]["color"] == "#8A8A8A"
    assert chart_element["categoryAxis"]["labelStyle"]["fontSize"] == 10.0
    assert chart_element["valueAxis"]["line"] == {
        "type": "solid", "color": "#888888", "width": 1.0,
    }
    assert chart_element["valueAxis"]["majorGridline"] == {
        "type": "solid", "color": "#3A3A3A", "width": 0.5,
    }
    assert chart_element["legend"]["position"] == "r"
    assert chart_element["legend"]["style"]["color"] == "#616878"
    assert chart_element["gapWidth"] == 175.0
    assert chart_element["overlap"] == 20.0
    assert chart_element["roundedCorners"] is True
    assert chart_element["plotArea"] == {
        "target": "outer", "x": 0.12, "y": 0.1, "w": 0.72, "h": 0.68,
    }
    labels = chart_element["series"][0]["dataLabels"]
    assert labels["showValue"] is True
    assert labels["position"] == "outEnd"
    assert labels["style"]["color"] == "#F2F2F2"

    spec = workspace / "update-chart.json"
    spec.write_text(json.dumps({
        "operations": [{
            "type": "update_chart",
            "slide": 1,
            "element_id": chart_element["elementId"],
            "title": "Updated sales",
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "East", "values": [12, 20, 28, 35]},
                {"name": "West", "values": [8, 15, 19, 26]},
            ],
            "show_legend": True,
            "series_colors": ["2F6B4F", "D97706"],
        }],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-chart-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=({
            "id": "chart-source",
            "name": "chart-source.pptx",
            "kind": "document",
            "local_path": str(source),
        },),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="update-chart.json",
            output_name="chart-updated.pptx",
            source_attachment_id="chart-source",
        )

    assert result.get("success") is True, result
    assert result["validation"]["chart_count"] == 1
    updated = Presentation(outputs / "chart-updated.pptx")
    updated_chart = updated.slides[0].shapes[0].chart
    assert updated_chart.chart_title.text_frame.text == "Updated sales"
    assert [str(label[0]) for label in updated_chart.plots[0].categories.flattened_labels] == [
        "Q1", "Q2", "Q3", "Q4",
    ]
    assert [series.name for series in updated_chart.series] == ["East", "West"]
    assert list(updated_chart.series[1].values) == [8.0, 15.0, 19.0, 26.0]
    assert updated_chart.has_legend is True


def test_presentation_chart_preserves_inherited_axis_and_gridline_defaults(tmp_path):
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "inherited-chart-axis.pptx"
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Sales", [10, 18, 24])
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(2), Cm(2), Cm(22), Cm(12),
        data,
    )
    deck.save(source)

    preview_dir = tmp_path / ".presentation" / "inherited-chart-axis"
    build_presentation_project(source, preview_dir)
    project = json.loads((preview_dir / "project.json").read_text(encoding="utf-8"))
    chart = next(
        element
        for element in project["slides"][0]["elements"]
        if element["elementType"] == "chart"
    )

    assert chart["categoryAxis"]["line"] == {
        "type": "solid", "color": "#888888", "width": 0.75,
    }
    assert chart["valueAxis"]["line"] == {
        "type": "solid", "color": "#888888", "width": 0.75,
    }
    assert chart["valueAxis"]["majorGridline"] == {
        "type": "solid", "color": "#D9DEE7", "width": 0.5,
    }


def test_presentation_preview_preserves_no_fill_and_no_line(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.xmlchemy import OxmlElement
    from xiaomei_brain.documents.presentation_project import (
        PROJECT_GENERATOR_VERSION,
        build_presentation_project,
    )

    source = tmp_path / "dark-cover.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x14)
    text_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(2), Cm(2), Cm(20), Cm(3),
    )
    text_shape.text = "把对话变成生产力"
    text_shape.fill.background()
    text_shape.line.fill.background()
    inherited_line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(2), Cm(6), Cm(20), Cm(3),
    )
    inherited_line_shape.text = "Theme line is not an explicit border"
    for line in inherited_line_shape._element.spPr.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ln"
    ):
        inherited_line_shape._element.spPr.remove(line)
    inherited_line_shape._element.spPr.append(OxmlElement("a:ln"))
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "dark-cover"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    element = project["slides"][0]["elements"][0]
    inherited_element = project["slides"][0]["elements"][1]

    assert project["generatorVersion"] == PROJECT_GENERATOR_VERSION
    assert project["slides"][0]["background"]["color"] == "#141414"
    assert element["fill"] == {"type": "none", "color": "transparent"}
    assert element["line"] == {
        "type": "none",
        "color": "transparent",
        "width": 0,
    }
    assert inherited_element["line"] == {
        "type": "none",
        "color": "transparent",
        "width": 0,
    }


def test_presentation_preview_preserves_line_width_markers_and_labels(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LABEL_POSITION, XL_MARKER_STYLE
    from pptx.util import Pt
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "styled-line.pptx"
    data = CategoryChartData()
    data.categories = ["Jan", "Feb", "Mar"]
    data.add_series("Users", [1, 5, 12])
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Cm(2), Cm(2), Cm(22), Cm(12),
        data,
    ).chart
    series = chart.series[0]
    series.format.line.color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    series.format.line.width = Pt(2)
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    series.marker.size = 8
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_value = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.ABOVE
    chart.plots[0].data_labels.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    chart.plots[0].data_labels.font.size = Pt(10)
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "styled-line"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    chart_element = project["slides"][0]["elements"][0]
    preview_series = chart_element["series"][0]

    assert preview_series["line"] == {
        "type": "solid", "color": "#C6F24E", "width": 2.0,
    }
    assert preview_series["marker"]["symbol"] == "circle"
    assert preview_series["marker"]["size"] == 8.0
    assert preview_series["marker"]["fill"] == {
        "type": "solid", "color": "#C6F24E",
    }
    assert preview_series["dataLabels"]["showValue"] is True
    assert preview_series["dataLabels"]["position"] == "t"
    assert preview_series["dataLabels"]["style"]["fontSize"] == 10.0


def test_presentation_preview_preserves_inherited_line_series(tmp_path):
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "inherited-line.pptx"
    data = CategoryChartData()
    data.categories = ["Jan", "Feb", "Mar"]
    data.add_series("Users", [1, 5, 12])
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Cm(2), Cm(2), Cm(22), Cm(12),
        data,
    )
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "inherited-line"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    chart_element = project["slides"][0]["elements"][0]
    preview_series = chart_element["series"][0]

    assert chart_element["chartType"] == "line"
    assert preview_series["line"] == {
        "type": "solid", "color": "#4F6BED", "width": 2.0,
    }
    assert preview_series["marker"] == {"symbol": "none", "size": 0}


def test_presentation_preview_extracts_scatter_and_radar_charts(tmp_path):
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "scatter-radar.pptx"
    deck = Presentation()

    scatter_slide = deck.slides.add_slide(deck.slide_layouts[6])
    scatter_data = XyChartData()
    scatter_series = scatter_data.add_series("Growth")
    for x_value, y_value in ((1, 4), (2, 9), (4, 16)):
        scatter_series.add_data_point(x_value, y_value)
    scatter_slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Cm(2), Cm(2), Cm(22), Cm(12),
        scatter_data,
    )

    radar_slide = deck.slides.add_slide(deck.slide_layouts[6])
    radar_data = CategoryChartData()
    radar_data.categories = ["Quality", "Speed", "Cost", "Service"]
    radar_data.add_series("Current", [82, 68, 74, 91])
    radar_slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR,
        Cm(2), Cm(2), Cm(22), Cm(12),
        radar_data,
    )
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "scatter-radar"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    scatter = project["slides"][0]["elements"][0]
    radar = project["slides"][1]["elements"][0]

    assert scatter["chartType"] == "xy_scatter"
    assert scatter["series"][0]["xValues"] == [1.0, 2.0, 4.0]
    assert scatter["series"][0]["values"] == [4.0, 9.0, 16.0]
    assert scatter["series"][0]["marker"] == {"symbol": "circle", "size": 5.0}
    assert scatter["series"][0]["line"]["type"] == "none"
    assert scatter["categoryAxis"]["position"] == "b"
    assert scatter["valueAxis"]["position"] == "l"

    assert radar["chartType"] == "radar"
    assert radar["categories"] == ["Quality", "Speed", "Cost", "Service"]
    assert radar["series"][0]["values"] == [82.0, 68.0, 74.0, 91.0]
    assert radar["series"][0]["line"]["type"] == "solid"


def test_presentation_preview_preserves_rich_text_layout_and_rotation(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "rich-cover.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(2), Cm(20), Cm(5))
    shape.rotation = 4
    shape.fill.background()
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Pt(12)
    frame.margin_right = Pt(14)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(8)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.space_after = Pt(5)
    first = paragraph.add_run()
    first.text = "Before "
    first.font.size = Pt(42)
    first.font.bold = True
    first.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    second = paragraph.add_run()
    second.text = "After"
    second.font.size = Pt(42)
    second.font.bold = True
    second.font.color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "rich-cover"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    element = project["slides"][0]["elements"][0]
    text_style = element["textStyle"]
    rich_paragraph = element["richText"]["paragraphs"][0]

    assert element["rotation"] == 4.0
    assert text_style["verticalAlign"] == "middle"
    assert text_style["margins"] == [6.0, 14.0, 8.0, 12.0]
    assert rich_paragraph["align"] == "center"
    assert rich_paragraph["spaceAfter"] == 5.0
    assert [run["color"] for run in rich_paragraph["runs"]] == ["#F2F2F2", "#C6F24E"]
    assert [run["fontSize"] for run in rich_paragraph["runs"]] == [42.0, 42.0]


def test_presentation_preview_preserves_table_cell_styles_and_merges(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "styled-table.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Cm(2), Cm(2), Cm(20), Cm(8)).table
    header = table.cell(0, 0)
    header.text = "Header"
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x14)
    header.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    header.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    header.merge(table.cell(0, 1))
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "styled-table"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    cells = project["slides"][0]["elements"][0]["cells"]

    assert cells[0]["fill"] == {"type": "solid", "color": "#141414"}
    assert cells[0]["textStyle"]["color"] == "#F2F2F2"
    assert cells[0]["textStyle"]["fontSize"] == 18.0
    assert cells[0]["columnSpan"] == 2
    assert cells[1]["hidden"] is True


def test_presentation_preview_preserves_group_transform_gradient_crop_and_shadow(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "visual-effects.pptx"
    image_path = tmp_path / "pixel.png"
    image_path.write_bytes(PNG_1PX)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])

    group = slide.shapes.add_group_shape()
    grouped = group.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(1), Cm(4), Cm(2))
    grouped.text = "Grouped"
    grouped.fill.gradient()
    grouped.fill.gradient_stops[0].color.rgb = RGBColor(0x14, 0x14, 0x14)
    grouped.fill.gradient_stops[-1].color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    grouped._element.spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}>'
        '<a:outerShdw blurRad="50800" dist="25400" dir="5400000">'
        '<a:srgbClr val="000000"><a:alpha val="40000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    ))
    group.left = Cm(3)
    group.top = Cm(4)
    group.width = Cm(12)
    group.height = Cm(6)

    picture = slide.shapes.add_picture(str(image_path), Cm(17), Cm(4), Cm(8), Cm(6))
    picture.crop_left = 0.1
    picture.crop_top = 0.2
    picture.crop_right = 0.15
    picture.crop_bottom = 0.05
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "visual-effects"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    elements = project["slides"][0]["elements"]
    shape = next(item for item in elements if item["elementType"] == "shape")
    image = next(item for item in elements if item["elementType"] == "image")

    assert shape["bounds"] == [
        round(float(group.left) / float(deck.slide_width) * 960, 3),
        round(float(group.top) / float(deck.slide_height) * project["size"][1], 3),
        round(float(group.width) / float(deck.slide_width) * 960, 3),
        round(float(group.height) / float(deck.slide_height) * project["size"][1], 3),
    ]
    assert shape["fill"]["type"] == "gradient"
    assert [stop["color"] for stop in shape["fill"]["stops"]] == ["#141414", "#C6F24E"]
    assert shape["shadow"] == {
        "blur": 4.0,
        "color": "#00000066",
        "offset": [0.0, 2.0],
    }
    assert image["fit"] == {"mode": "fill"}
    assert image["crop"] == {
        "left": 0.1, "top": 0.2, "right": 0.15, "bottom": 0.05,
    }


def test_presentation_preview_extracts_connectors_arrows_and_theme_lines(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn
    from pptx.util import Pt
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "connectors.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    straight = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Cm(2), Cm(2), Cm(12), Cm(5),
    )
    straight.line.color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    straight.line.width = Pt(2)
    line = straight._element.spPr.find(qn("a:ln"))
    line.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="dash"/>'))
    line.append(parse_xml(
        f'<a:headEnd {nsdecls("a")} type="diamond" w="sm" len="med"/>'
    ))
    line.append(parse_xml(
        f'<a:tailEnd {nsdecls("a")} type="triangle" w="lg" len="lg"/>'
    ))
    slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, Cm(2), Cm(8), Cm(12), Cm(12),
    )
    slide.shapes.add_connector(
        MSO_CONNECTOR.CURVE, Cm(16), Cm(2), Cm(26), Cm(12),
    )
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "connectors"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    lines = [
        element for element in project["slides"][0]["elements"]
        if element["elementType"] == "line"
    ]

    assert len(lines) == 3
    assert [item["connectorKind"] for item in lines] == [
        "line", "bentConnector3", "curvedConnector3",
    ]
    assert lines[0]["line"] == {
        "type": "dash", "color": "#C6F24E", "width": 2.0,
    }
    assert lines[0]["startArrow"] == {
        "type": "diamond", "width": "sm", "length": "med",
    }
    assert lines[0]["endArrow"] == {
        "type": "triangle", "width": "lg", "length": "lg",
    }
    assert lines[1]["line"]["color"] != "transparent"
    assert lines[1]["line"]["width"] > 0

    extracted = PresentationExtractor().extract(source)
    index_text = extracted.sections[0].content
    assert "[元素索引]" in index_text
    assert f'element_id="slide-1-shape-id-{straight.shape_id}" type=line' in index_text
    assert "position_cm=" in index_text
    assert extracted.sections[0].metadata["element_count"] == 3


def test_write_document_updates_connector_and_freeform_styles(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    workspace.mkdir()
    source = tmp_path / "editable-shapes.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_textbox(Cm(1), Cm(1), Cm(8), Cm(1)).text = "Process"
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, Cm(2), Cm(4), Cm(14), Cm(8),
    )
    builder = slide.shapes.build_freeform(0, 0, scale=(10000, 10000))
    builder.add_line_segments([(80, 0), (100, 80), (50, 50), (0, 80)], close=True)
    freeform = builder.convert_to_shape(Cm(17), Cm(4))
    freeform.fill.solid()
    freeform.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    deck.save(source)

    preview_dir = tmp_path / ".presentation" / "editable-shapes"
    build_presentation_project(source, preview_dir)
    project = json.loads((preview_dir / "project.json").read_text(encoding="utf-8"))
    elements = project["slides"][0]["elements"]
    connector_id = next(item["elementId"] for item in elements if item["elementType"] == "line")
    freeform_id = next(
        item["elementId"] for item in elements
        if item["elementType"] == "shape" and item.get("customGeometry")
    )
    spec = workspace / "shape-update.json"
    spec.write_text(json.dumps({
        "operations": [
            {
                "type": "update_element",
                "slide": 1,
                "element_id": connector_id,
                "line_color": "2F6B4F",
                "line_width_pt": 3,
                "line_dash": "dash_dot",
                "start_arrow": "diamond",
                "end_arrow": {"type": "triangle", "width": "lg", "length": "lg"},
                "line_transparency": 25,
            },
            {
                "type": "update_element",
                "slide": 1,
                "element_id": freeform_id,
                "fill_color": "C6F24E",
                "fill_transparency": 40,
                "line_color": "141414",
            },
        ],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-shape-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=({
            "id": "source-deck",
            "name": source.name,
            "kind": "document",
            "local_path": str(source),
        },),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="shape-update.json",
            output_name="updated-shapes.pptx",
            source_attachment_id="source-deck",
        )

    assert result.get("success") is True, result
    updated = Presentation(outputs / "updated-shapes.pptx")
    updated_connector = next(
        shape for shape in updated.slides[0].shapes if shape.shape_id == connector.shape_id
    )
    line = updated_connector._element.spPr.find(qn("a:ln"))
    assert line.find(qn("a:prstDash")).get("val") == "dashDot"
    assert line.find(qn("a:headEnd")).get("type") == "diamond"
    assert line.find(qn("a:tailEnd")).get("type") == "triangle"
    assert line.find(qn("a:solidFill"))[0].find(qn("a:alpha")).get("val") == "75000"
    updated_freeform = next(
        shape for shape in updated.slides[0].shapes if shape.shape_id == freeform.shape_id
    )
    fill = updated_freeform._element.spPr.find(qn("a:solidFill"))
    assert fill[0].find(qn("a:alpha")).get("val") == "60000"


def test_presentation_preview_extracts_freeform_geometry(tmp_path):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "freeform.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    builder = slide.shapes.build_freeform(0, 0, scale=(10000, 10000))
    builder.add_line_segments([(80, 0), (100, 80), (50, 50), (0, 80)], close=True)
    shape = builder.convert_to_shape(Cm(4), Cm(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xC6, 0xF2, 0x4E)
    shape.line.color.rgb = RGBColor(0x14, 0x14, 0x14)
    shape.line.width = Pt(1.5)
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "freeform"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    element = project["slides"][0]["elements"][0]
    paths = element["customGeometry"]["paths"]

    assert element["elementType"] == "shape"
    assert element["fill"] == {"type": "solid", "color": "#C6F24E"}
    assert paths
    assert paths[0]["d"].startswith("M ")
    assert " L " in paths[0]["d"]
    assert paths[0]["d"].endswith("Z")


def test_presentation_preview_extracts_native_formula_and_embedded_video(tmp_path):
    from lxml import etree
    from pptx.oxml.ns import qn
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "formula-and-video.pptx"
    video_path = tmp_path / "sample.mp4"
    audio_path = tmp_path / "sample.wav"
    poster_path = tmp_path / "poster.png"
    video_path.write_bytes(b"fake-mp4-for-package-test")
    audio_path.write_bytes(b"fake-wav-for-package-test")
    poster_path.write_bytes(PNG_1PX)

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    formula_shape = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(12), Cm(3))
    paragraph = formula_shape._element.find(".//" + qn("a:p"))
    assert paragraph is not None
    for child in list(paragraph):
        paragraph.remove(child)
    formula = etree.fromstring("""
      <a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main"
             xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">
        <m:oMath>
          <m:f>
            <m:num><m:r><m:t>x</m:t></m:r></m:num>
            <m:den>
              <m:sSup>
                <m:e><m:r><m:t>y</m:t></m:r></m:e>
                <m:sup><m:r><m:t>2</m:t></m:r></m:sup>
              </m:sSup>
            </m:den>
          </m:f>
        </m:oMath>
      </a14:m>
    """.strip().encode("utf-8"))
    paragraph.append(formula)
    slide.shapes.add_movie(
        str(video_path),
        Cm(3), Cm(6), Cm(16), Cm(9),
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )
    audio_shape = slide.shapes.add_movie(
        str(audio_path),
        Cm(21), Cm(6), Cm(10), Cm(3),
        poster_frame_image=str(poster_path),
        mime_type="audio/wav",
    )
    audio_node = audio_shape._element.find(".//" + qn("a:videoFile"))
    assert audio_node is not None
    audio_node.tag = qn("a:audioFile")
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "formula-and-video"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    elements = project["slides"][0]["elements"]
    formula_element = next(item for item in elements if item["elementType"] == "formula")
    assert "<mfrac>" in formula_element["mathMl"]
    assert "<msup>" in formula_element["mathMl"]
    assert formula_element["fallbackText"] == "xy2"
    media_element = next(
        item for item in elements
        if item["elementType"] == "media" and item["mediaKind"] == "video"
    )
    assert media_element["mediaKind"] == "video"
    assert media_element["mimeType"] == "video/mp4"
    assert (project_dir / media_element["src"]).read_bytes() == video_path.read_bytes()
    assert (project_dir / media_element["posterSrc"]).is_file()
    audio_element = next(
        item for item in elements
        if item["elementType"] == "media" and item["mediaKind"] == "audio"
    )
    assert audio_element["mimeType"].startswith("audio/")
    assert (project_dir / audio_element["src"]).read_bytes() == audio_path.read_bytes()


def test_write_document_creates_and_updates_native_formula_and_media(tmp_path):
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    registry = _presentation_registry()
    tool = create_write_document_tool(registry)
    workspace = tmp_path / "workspace"
    outputs = workspace / "outputs"
    work = workspace / "work"
    work.mkdir(parents=True)
    (work / "poster.png").write_bytes(PNG_1PX)
    (work / "clip.mp4").write_bytes(b"first-video")
    (work / "replacement.mp4").write_bytes(b"replacement-video")
    spec = workspace / "formula-media.json"
    spec.write_text(json.dumps({
        "slides": [{
            "type": "blank",
            "elements": [
                {
                    "type": "formula",
                    "name": "NativeFormula",
                    "x_cm": 2, "y_cm": 1, "width_cm": 14, "height_cm": 3,
                    "expression": {
                        "type": "fraction",
                        "numerator": "x",
                        "denominator": {
                            "type": "superscript",
                            "base": "y",
                            "superscript": 2,
                        },
                    },
                },
                {
                    "type": "media",
                    "media_kind": "video",
                    "workspace_path": "work/clip.mp4",
                    "poster_workspace_path": "work/poster.png",
                    "x_cm": 3, "y_cm": 5, "width_cm": 16, "height_cm": 9,
                },
            ],
        }],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-formula-media",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=(),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        result = tool.execute(
            format="presentation",
            specification_path="formula-media.json",
            output_name="formula-media.pptx",
        )

    assert result.get("success") is True, result
    assert result["validation"]["formula_count"] == 1
    assert result["validation"]["media_count"] == 1
    project_dir = tmp_path / ".presentation" / "formula-media-created"
    build_presentation_project(outputs / "formula-media.pptx", project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    elements = project["slides"][0]["elements"]
    formula = next(item for item in elements if item["elementType"] == "formula")
    media = next(item for item in elements if item["elementType"] == "media")
    assert "<mfrac>" in formula["mathMl"]
    assert "<msup>" in formula["mathMl"]
    assert (project_dir / media["src"]).read_bytes() == b"first-video"

    update = workspace / "formula-media-update.json"
    update.write_text(json.dumps({
        "operations": [
            {
                "type": "update_formula",
                "slide": 1,
                "element_id": formula["elementId"],
                "expression": {
                    "type": "radical",
                    "degree": 3,
                    "radicand": "z",
                },
            },
            {
                "type": "replace_media",
                "slide": 1,
                "element_id": media["elementId"],
                "media_kind": "video",
                "workspace_path": "work/replacement.mp4",
                "poster_workspace_path": "work/poster.png",
            },
        ],
    }), encoding="utf-8")

    with bind_tool_execution(
        tool_call_id="call-presentation-formula-media-update",
        tool_name="write_document",
        arguments={},
        artifact_callback=None,
        attachments=({
            "id": "source-deck",
            "name": "formula-media.pptx",
            "kind": "document",
            "local_path": str(outputs / "formula-media.pptx"),
        },),
        workspace_root=str(workspace),
        output_root=str(outputs),
    ):
        updated = tool.execute(
            format="presentation",
            specification_path="formula-media-update.json",
            output_name="formula-media-updated.pptx",
            source_attachment_id="source-deck",
        )

    assert updated.get("success") is True, updated
    updated_dir = tmp_path / ".presentation" / "formula-media-updated"
    build_presentation_project(outputs / "formula-media-updated.pptx", updated_dir)
    updated_project = json.loads(
        (updated_dir / "project.json").read_text(encoding="utf-8")
    )
    updated_elements = updated_project["slides"][0]["elements"]
    updated_formula = next(
        item for item in updated_elements if item["elementType"] == "formula"
    )
    updated_media = next(
        item for item in updated_elements if item["elementType"] == "media"
    )
    assert "<mroot>" in updated_formula["mathMl"]
    assert (updated_dir / updated_media["src"]).read_bytes() == b"replacement-video"


def test_presentation_preview_extracts_transition_and_animation(tmp_path):
    from lxml import etree
    from xiaomei_brain.documents.presentation_project import build_presentation_project

    source = tmp_path / "animated.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_textbox(Cm(2), Cm(2), Cm(10), Cm(2))
    shape.text = "Animated title"
    transition = etree.fromstring("""
      <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    spd="slow" advClick="1" advTm="3200">
        <p:fade thruBlk="0"/>
      </p:transition>
    """.strip().encode("utf-8"))
    timing = etree.fromstring(f"""
      <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:tnLst><p:par><p:cTn id="1" nodeType="tmRoot"><p:childTnLst>
          <p:seq><p:cTn id="2" nodeType="clickEffect"><p:childTnLst>
            <p:animEffect transition="in" filter="fade">
              <p:cBhvr>
                <p:cTn id="3" dur="650">
                  <p:stCondLst><p:cond delay="120"/></p:stCondLst>
                </p:cTn>
                <p:tgtEl><p:spTgt spid="{shape.shape_id}"/></p:tgtEl>
              </p:cBhvr>
            </p:animEffect>
          </p:childTnLst></p:cTn></p:seq>
        </p:childTnLst></p:cTn></p:par></p:tnLst>
      </p:timing>
    """.strip().encode("utf-8"))
    slide._element.append(transition)
    slide._element.append(timing)
    deck.save(source)

    project_dir = tmp_path / ".presentation" / "animated"
    build_presentation_project(source, project_dir)
    project = json.loads((project_dir / "project.json").read_text(encoding="utf-8"))
    preview_slide = project["slides"][0]
    assert preview_slide["transition"] == {
        "type": "fade",
        "durationMs": 1000,
        "advanceOnClick": True,
        "advanceAfterMs": 3200,
        "thruBlk": "0",
    }
    animation = preview_slide["animations"][0]
    assert animation["targetShapeId"] == shape.shape_id
    assert animation["targetElementId"] == f"slide-1-shape-id-{shape.shape_id}"
    assert animation["effect"] == "fade"
    assert animation["durationMs"] == 650
    assert animation["delayMs"] == 120
    assert preview_slide["elements"][0]["animations"][0]["id"] == animation["id"]


def test_write_document_creates_and_updates_basic_transition_and_animation(tmp_path):
    from xiaomei_brain.plugins.tools.document_presentation.writer import PresentationWriter

    writer = PresentationWriter()
    created = tmp_path / "motion-created.pptx"
    writer.write({
        "slides": [{
            "type": "blank",
            "transition": {
                "type": "push",
                "speed": "slow",
                "direction": "right",
                "advance_on_click": True,
            },
            "elements": [{
                "type": "text",
                "text": "关键结论",
                "x_cm": 3, "y_cm": 4, "width_cm": 16, "height_cm": 3,
                "animation": {
                    "effect": "fade",
                    "trigger": "after_previous",
                    "duration_ms": 800,
                    "delay_ms": 150,
                },
            }],
        }],
    }, created)

    created_project = json.loads(
        (tmp_path / ".presentation" / "motion-created" / "project.json").read_text(
            encoding="utf-8"
        )
    )
    created_slide = created_project["slides"][0]
    assert created_slide["transition"] == {
        "type": "push",
        "durationMs": 1000,
        "advanceOnClick": True,
        "dir": "r",
    }
    first_animation = created_slide["animations"][0]
    assert first_animation["effect"] == "fade"
    assert first_animation["trigger"] == "afterEffect"
    assert first_animation["durationMs"] == 800
    assert first_animation["delayMs"] == 150
    element_id = created_slide["elements"][0]["elementId"]

    updated = tmp_path / "motion-updated.pptx"
    writer.write({
        "operations": [
            {
                "type": "set_transition",
                "slide": 1,
                "transition": {
                    "type": "wipe",
                    "speed": "fast",
                    "direction": "down",
                    "advance_on_click": False,
                    "advance_after_ms": 2500,
                },
            },
            {
                "type": "add_animation",
                "slide": 1,
                "element_id": element_id,
                "animation": {
                    "effect": "fly",
                    "direction": "left",
                    "trigger": "on_click",
                    "duration_ms": 500,
                },
            },
        ],
    }, updated, source_path=created)

    updated_project = json.loads(
        (tmp_path / ".presentation" / "motion-updated" / "project.json").read_text(
            encoding="utf-8"
        )
    )
    updated_slide = updated_project["slides"][0]
    assert updated_slide["transition"] == {
        "type": "wipe",
        "durationMs": 250,
        "advanceOnClick": False,
        "advanceAfterMs": 2500,
        "dir": "d",
    }
    assert len(updated_slide["animations"]) == 2
    second_animation = updated_slide["animations"][1]
    assert second_animation["effect"] == "fly(fromLeft)"
    assert second_animation["trigger"] == "clickEffect"
    assert second_animation["durationMs"] == 500


def test_presentation_preview_reads_smartart_data_model():
    from lxml import etree
    from xiaomei_brain.documents.presentation_project import _smartart_element

    class RelatedPart:
        def __init__(self, blob: bytes):
            self.blob = blob

    class ShapePart:
        def __init__(self, related):
            self._related = related

        def related_part(self, relationship_id):
            return self._related[relationship_id]

    class Shape:
        pass

    shape = Shape()
    shape._element = etree.fromstring("""
      <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                      xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
                      xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
        <a:graphic><a:graphicData><dgm:relIds r:dm="rIdData" r:lo="rIdLayout"/></a:graphicData></a:graphic>
      </p:graphicFrame>
    """.strip().encode("utf-8"))
    shape.part = ShapePart({
        "rIdData": RelatedPart("""
          <dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <dgm:ptLst>
              <dgm:pt modelId="one"><dgm:t><a:p><a:r><a:t>需求</a:t></a:r></a:p></dgm:t></dgm:pt>
              <dgm:pt modelId="two"><dgm:t><a:p><a:r><a:t>交付</a:t></a:r></a:p></dgm:t></dgm:pt>
            </dgm:ptLst>
            <dgm:cxnLst><dgm:cxn srcId="one" destId="two" type="parOf"/></dgm:cxnLst>
          </dgm:dataModel>
        """.strip().encode("utf-8")),
        "rIdLayout": RelatedPart("""
          <dgm:layoutDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">
            <dgm:title val="Basic Process"/>
          </dgm:layoutDef>
        """.strip().encode("utf-8")),
    })

    element = _smartart_element(shape, "slide-1-shape-id-7", [10, 20, 600, 240])

    assert element is not None
    assert element["elementType"] == "smartart"
    assert element["layoutName"] == "Basic Process"
    assert [node["text"] for node in element["nodes"]] == ["需求", "交付"]
    assert element["connections"] == [{
        "source": "one", "target": "two", "type": "parOf",
    }]
